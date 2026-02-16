from __future__ import annotations

import os
import re
import json
import sys
import tempfile
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable, Iterator, List, Optional, Tuple
from urllib.parse import urlparse, unquote
import fnmatch

import requests
from serial.tools import list_ports
import shutil


# =====================
# 开发板预设（PlatformIO board + platform）
# =====================
# 每条: (显示名称, platform, board_id)
SUPPORTED_BOARDS: List[Tuple[str, str, str]] = [
    ("ESP8266 NodeMCU (D1 mini 等)", "espressif8266", "nodemcuv2"),
    ("ESP32 Dev Module", "espressif32", "esp32dev"),
    ("ESP32-S3 DevKitC", "espressif32", "esp32-s3-devkitc-1"),
    ("ESP32-C3 DevModule", "espressif32", "esp32-c3-devkitm-1"),
    ("树莓派 Pico (RP2040)", "raspberrypi", "pico"),
    ("STM32H743 通用", "ststm32", "genericSTM32H743VITx"),
    ("STM32F401 Black Pill", "ststm32", "blackpill_f401cc"),
    ("Arduino Uno", "atmelavr", "uno"),
]


def get_supported_boards() -> List[dict]:
    """返回供前端下拉使用的开发板列表，每项含 id（board_id）、platform、name"""
    return [
        {"id": board_id, "platform": platform, "name": name}
        for name, platform, board_id in SUPPORTED_BOARDS
    ]


# =====================
# 1. USB / 串口 设备感知
# =====================

_DEVICES_CACHE_TTL = 2.0  # 秒
_devices_cache: Optional[tuple] = None  # (timestamp, result)


def list_serial_devices(force_refresh: bool = False) -> List[dict]:
    """列出当前所有串口设备，返回简单信息列表。2 秒内重复调用返回缓存，减少 USB 扫描。force_refresh=True 时跳过缓存。"""
    global _devices_cache
    now = time.monotonic()
    if not force_refresh and _devices_cache is not None:
        ts, cached = _devices_cache
        if now - ts < _DEVICES_CACHE_TTL:
            return cached
    devices: List[dict] = []
    for p in list_ports.comports():
        devices.append(
            {
                "device": p.device,  # 如 /dev/cu.usbserial-0001
                "description": p.description,  # 设备名称
                "hwid": p.hwid,  # VID/PID 等
                "manufacturer": getattr(p, "manufacturer", None),
                "product": getattr(p, "product", None),
            }
        )

    # 过滤掉明显不是开发板的串口（如 Mac 的 debug-console、蓝牙虚拟串口）
    bad_keywords = ["debug-console", "bluetooth"]
    filtered: List[dict] = []
    for d in devices:
        text = (
            (d.get("device") or "")
            + " "
            + (d.get("description") or "")
            + " "
            + (d.get("product") or "")
        ).lower()
        if any(bad in text for bad in bad_keywords):
            continue
        filtered.append(d)

    # 如果全被筛掉，就退回原列表，防止用户完全看不到任何设备
    result = filtered or devices
    _devices_cache = (time.monotonic(), result)
    return result


def guess_esp8266_port(devices: List[dict]) -> Optional[str]:
    """根据常见字段粗略猜测哪个串口是 ESP8266"""
    keywords = ["CP210", "CH340", "ESP8266", "USB-SERIAL", "USB2.0-Serial"]
    for d in devices:
        desc = (d.get("description") or "") + " " + (d.get("product") or "")
        if any(k.lower() in desc.lower() for k in keywords):
            return d["device"]
    # 如果猜不到，就返回第一个串口
    return devices[0]["device"] if devices else None


# =====================
# 2. 调用 Qwen Coder 2.5 生成代码
# =====================

def extract_code_from_md(text: str) -> str:
    """从可能带 ``` 或 ```python 的内容里提取纯代码"""
    code_blocks = re.findall(r"```(?:[a-zA-Z0-9_+-]+)?\s*([\\s\\S]*?)```", text)
    if code_blocks:
        return code_blocks[0].strip()
    return text.strip()


# 自动选择连接成功的 API：缓存当前可用的 endpoint，避免每次按固定优先级
_working_endpoint_cache: Optional[dict] = None  # {"expiry": float, "url": str, "headers": dict, "provider": str}
_WORKING_ENDPOINT_TTL = 60  # 缓存 60 秒后重新探测


def _build_endpoint_candidates() -> List[tuple]:
    """
    按当前环境变量构造所有已配置的 (url, headers, provider_key) 候选。
    默认顺序：本地 Qwen > DeepSeek 联网 > DashScope。
    若设置 PREFER_DEEPSEEK=1 或 DEEPSEEK_PREFER=1，则优先尝试 DeepSeek（便于只用 DeepSeek 时避免先连本地 Qwen 超时）。
    """
    candidates: List[tuple] = []
    base = os.environ.get("QWEN_API_BASE")
    api_key = os.environ.get("QWEN_API_KEY")
    if base:
        base = base.rstrip("/")
        url = f"{base}/chat/completions" if not base.endswith("/chat/completions") else base
        headers = {"Content-Type": "application/json"}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        candidates.append((url, headers, "qwen_local"))

    deepseek_key = os.environ.get("DEEPSEEK_API_KEY")
    if deepseek_key:
        deepseek_base = (os.environ.get("DEEPSEEK_API_BASE") or "https://api.deepseek.com").rstrip("/")
        url = f"{deepseek_base}/v1/chat/completions"
        headers = {"Authorization": f"Bearer {deepseek_key}", "Content-Type": "application/json"}
        candidates.append((url, headers, "deepseek"))

    dash_api_key = os.environ.get("DASHSCOPE_API_KEY")
    if dash_api_key:
        url = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
        headers = {"Authorization": f"Bearer {dash_api_key}", "Content-Type": "application/json"}
        candidates.append((url, headers, "dashscope"))

    prefer_deepseek = os.environ.get("PREFER_DEEPSEEK", "") or os.environ.get("DEEPSEEK_PREFER", "")
    if prefer_deepseek and str(prefer_deepseek).lower() in ("1", "true", "yes"):
        # 把 DeepSeek 移到最前，避免本地 Qwen 未启动时长时间等待
        deepseek_items = [(u, h, p) for u, h, p in candidates if p == "deepseek"]
        others = [(u, h, p) for u, h, p in candidates if p != "deepseek"]
        candidates = deepseek_items + others

    return candidates


def _try_endpoint(url: str, headers: dict, provider: str, timeout: int = 8) -> Tuple[bool, Optional[str]]:
    """
    对给定 endpoint 发一次最小 ping 请求。
    返回 (是否成功, 失败时的简短原因，成功时为 None)。
    按 provider 使用对应模型名，避免向本地 Qwen 误发 deepseek-chat。
    """
    if provider == "deepseek":
        model = os.environ.get("DEEPSEEK_MODEL", "deepseek-chat")
    else:
        model = os.environ.get("LUMI_MODEL", "qwen2.5-coder-14b")
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": "ping"}],
        "max_tokens": 1,
        "temperature": 0,
    }
    try:
        r = requests.post(url, headers=headers, json=payload, timeout=timeout)
        if r.status_code != 200:
            try:
                body = r.json()
                msg = body.get("error", {}).get("message", body.get("message", "")) or r.text[:80]
            except Exception:
                msg = r.text[:80] if r.text else f"HTTP {r.status_code}"
            msg = (msg or f"HTTP {r.status_code}").strip()
            if "insufficient balance" in msg.lower() or "余额不足" in msg:
                msg = "DeepSeek 账户余额不足，请前往控制台充值或查看计费"
            return False, msg
        r.json()
        return True, None
    except requests.exceptions.Timeout:
        return False, "连接超时"
    except requests.exceptions.ConnectionError:
        return False, "连接被拒绝或网络不可达"
    except Exception as e:
        err = str(e)[:120].strip()
        if "insufficient balance" in err.lower():
            err = "DeepSeek 账户余额不足，请前往控制台充值或查看计费"
        return False, err


def _get_working_endpoint() -> tuple[str, dict]:
    """
    从已配置的候选中自动选择第一个连接成功的 API，并短期缓存。
    若均未配置或均连接失败则抛出 RuntimeError（含各接口失败原因，便于排查）。
    """
    global _working_endpoint_cache
    now = time.monotonic()
    if _working_endpoint_cache and now < _working_endpoint_cache.get("expiry", 0):
        return _working_endpoint_cache["url"], _working_endpoint_cache["headers"]

    candidates = _build_endpoint_candidates()
    if not candidates:
        raise RuntimeError(
            "未配置模型接口：请设置 QWEN_API_BASE（本地）、DEEPSEEK_API_KEY（DeepSeek 联网）或 DASHSCOPE_API_KEY 之一。"
        )

    errors = []
    for url, headers, provider in candidates:
        ok, err_msg = _try_endpoint(url, headers, provider)
        if ok:
            _working_endpoint_cache = {
                "expiry": now + _WORKING_ENDPOINT_TTL,
                "url": url,
                "headers": headers,
                "provider": provider,
            }
            return url, headers
        errors.append(f"{provider}: {err_msg or '失败'}")
    raise RuntimeError(
        "所有已配置的模型接口均无法连接。已尝试: " + "；".join(errors)
    )


def _get_qwen_endpoint() -> tuple[str, dict]:
    """
    统一返回当前可用的模型接口 URL 和 headers（OpenAI 兼容格式）。
    自动选择连接成功的 API，并缓存约 60 秒。
    """
    return _get_working_endpoint()


def _get_effective_model(default: str) -> str:
    """根据当前选中的 API（或环境变量）返回对应模型名。"""
    if _working_endpoint_cache and time.monotonic() < _working_endpoint_cache.get("expiry", 0):
        p = _working_endpoint_cache.get("provider", "")
        if p == "deepseek":
            return os.environ.get("DEEPSEEK_MODEL", "deepseek-chat")
        if p in ("qwen_local", "dashscope"):
            return os.environ.get("LUMI_MODEL", default)
    if os.environ.get("DEEPSEEK_API_KEY"):
        return os.environ.get("DEEPSEEK_MODEL", "deepseek-chat")
    return os.environ.get("LUMI_MODEL", default)


def get_model_provider_info() -> dict:
    """
    返回当前使用的模型接口类型，供开发者调试展示。
    若已通过自动选择得到可用 endpoint，则返回该接口；否则按配置顺序返回第一个已配置的（可能尚未探测）。
    展示的模型名与接口一致：DeepSeek 用 DEEPSEEK_MODEL，其余用 LUMI_MODEL。
    """
    labels = {
        "qwen_local": "您的本地API",
        "deepseek": "DeepSeek 联网 API",
        "dashscope": "DashScope（阿里云）",
    }
    # 按接口类型取模型名，避免「显示 Qwen 接口却显示 deepseek-chat」的不一致
    models = {
        "qwen_local": os.environ.get("LUMI_MODEL", "qwen2.5-coder-14b"),
        "deepseek": os.environ.get("DEEPSEEK_MODEL", "deepseek-chat"),
        "dashscope": os.environ.get("LUMI_MODEL", "qwen2.5-coder-14b"),
    }
    if _working_endpoint_cache and time.monotonic() < _working_endpoint_cache.get("expiry", 0):
        p = _working_endpoint_cache.get("provider", "")
        return {
            "provider": p,
            "label": labels.get(p, p) + "（已连接）",
            "model": models.get(p, ""),
        }
    candidates = _build_endpoint_candidates()
    if not candidates:
        return {"provider": "none", "label": "未配置", "model": ""}
    _, __, p = candidates[0]
    return {
        "provider": p,
        "label": labels.get(p, p),
        "model": models.get(p, ""),
    }


def _post_chat_with_retry(
    url: str, headers: dict, payload: dict, *, retries: int = 1, timeout: int = 60
):
    """
    带重试和更友好错误信息的请求封装。
    - 对超时 / 连接错误会自动重试
    - 对 4xx / 5xx 返回更清晰的提示
    """
    last_err: Exception | None = None
    for attempt in range(retries + 1):
        try:
            resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
            resp.raise_for_status()
            return resp
        except requests.exceptions.Timeout as e:
            last_err = e
            if attempt == retries:
                raise RuntimeError(
                    "模型接口请求超时，请检查本地服务是否繁忙或无响应。"
                    " 若内容较长可设置环境变量 QWEN_REQUEST_TIMEOUT=600（秒）后重试。"
                ) from e
        except requests.exceptions.ConnectionError as e:
            last_err = e
            if attempt == retries:
                raise RuntimeError(
                    "无法连接到本地模型接口，请确认 QWEN_API_BASE 地址正确且服务已启动。"
                ) from e
        except requests.exceptions.HTTPError as e:
            # 对 4xx/5xx 直接给出状态码与返回体前一部分
            last_err = e
            status = e.response.status_code if e.response is not None else "unknown"
            text = ""
            if e.response is not None:
                try:
                    text = e.response.text[:300]
                except Exception:
                    text = "<response decoding error>"
            raise RuntimeError(f"模型接口返回 HTTP {status} 错误：{text}") from e
        except Exception as e:  # 意料之外的异常
            last_err = e
            if attempt == retries:
                raise
    if last_err:
        raise last_err
    raise RuntimeError("未知错误：模型请求失败。")


def call_qwen_coder(user_instruction: str, model: str = "qwen2.5-coder-14b") -> str:
    """
    调用 Qwen Coder 2.5，根据用户需求生成 MicroPython 脚本（运行在 ESP8266 上）。
    优先使用本地 / 自建服务（QWEN_API_BASE），否则回退到 DashScope。
    """
    url, headers = _get_qwen_endpoint()

    system_prompt = (
        "You are a senior embedded / MicroPython engineer.\\n"
        "Target platform: ESP8266 running MicroPython firmware.\\n"
        "Requirements:\\n"
        "1. Only output a complete script that can be saved as main.py.\\n"
        "2. Do NOT output explanations, Markdown, or ``` fences.\\n"
        "3. For debug prints, just use print().\\n"
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_instruction},
    ]

    payload = {
        "model": _get_effective_model(model),
        "messages": messages,
        "temperature": 0.2,
    }

    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=180)
    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    return extract_code_from_md(content)


def call_qwen_cpp_for_platformio(
    user_instruction: str, model: str = "qwen2.5-coder-14b"
) -> str:
    """
    调用 Qwen Coder 2.5 生成可由 PlatformIO 编译的 C++ 主文件（src/main.cpp）。

    设计目标：
    - 支持 ESP8266 / ESP32 / STM32（包括 H743 等）等常见单片机；
    - 使用 Arduino 风格的 setup()/loop()，由 PlatformIO 选择具体平台与核心；
    - 只返回 main.cpp 的完整代码，不包含 platformio.ini 或其它文件。
    """
    url, headers = _get_qwen_endpoint()

    system_prompt = (
        "You are a senior embedded C++ engineer.\n"
        "Target: Arduino-style C++ project compiled by PlatformIO for boards such as "
        "ESP8266, ESP32, Raspberry Pi Pico (RP2040), or STM32 (including STM32H743 based quadcopter flight controllers).\n"
        "Requirements:\n"
        "1. Only output a single complete C++ source file that can be saved as src/main.cpp.\n"
        "2. Use Arduino-style setup()/loop() and APIs when possible so that the same code "
        "can run on ESP, STM32, or Raspberry Pi Pico (via Arduino core). For ESP8266 compatibility, define "
        "entry points as: extern \"C\" void setup() { ... } and extern \"C\" void loop() { ... }.\n"
        "2a. On ESP8266 do NOT use C++ STL (no std::vector, no std::initializer_list, no #include <vector>); "
        "use only #include <Arduino.h>, Arduino APIs, and C-style arrays.\n"
        "2b. On Raspberry Pi Pico (RP2040) use Arduino-Pico core APIs; GPIO numbers are 0–29, onboard LED usually GPIO 25.\n"
        "3. Do NOT output platformio.ini or any other files.\n"
        "4. Do NOT output Markdown, explanations, or ``` fences. Comments inside code are OK.\n"
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_instruction},
    ]

    payload = {
        "model": _get_effective_model(model),
        "messages": messages,
        "temperature": 0.2,
    }

    resp = _post_chat_with_retry(
        _get_qwen_endpoint()[0],
        _get_qwen_endpoint()[1],
        payload,
        retries=1,
        timeout=180,
    )
    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    return extract_code_from_md(content)


def call_qwen_file_editor(
    file_path: str,
    file_content: str,
    user_instruction: str,
    selected_text: Optional[str] = None,
    context_files: Optional[List[dict]] = None,
    model: str = "qwen2.5-coder-14b",
) -> str:
    """
    使用 Qwen Coder 2.5 对指定文件内容进行编辑，返回完整的新文件内容。
    - selected_text: 用户选中的代码片段，作为“当前焦点”上下文，便于精准修改。
    - context_files: 引用文件列表 [{"path": "rel/path", "content": "..."}]，用于跨文件理解。
    """
    url, headers = _get_qwen_endpoint()

    system_prompt = (
        "You are a senior AI pair-programming assistant. You modify code/files based on natural language instructions, "
        "with full context of the current file and optional selected snippet or referenced files.\n"
        "Requirements:\n"
        "1. You will see the current file path, its full content, and the user's editing request.\n"
        "2. If a 'selected snippet' is provided, that is the focus area the user is referring to.\n"
        "3. If 'referenced files' are provided, use them for context (e.g. imports, types) but only output the NEW content for the MAIN file.\n"
        "4. Output ONLY the FULL new file content for the main file. No explanations, no Markdown fences.\n"
        "5. Preserve file structure and style; only change what the instruction asks for.\n"
    )

    user_parts = [
        f"文件路径: {file_path}",
        f"用户的修改需求：\n{user_instruction}",
        "",
        "当前文件完整内容：",
        "```text",
        file_content,
        "```",
    ]
    if selected_text and selected_text.strip():
        user_parts.extend(["", "用户选中的代码（当前焦点）：", "```", selected_text.strip(), "```"])
    if context_files:
        user_parts.append("")
        user_parts.append("引用文件（仅作上下文，不要输出这些文件）：")
        for ctx in context_files[:10]:
            path = ctx.get("path") or "?"
            content = (ctx.get("content") or "")[:8000]
            user_parts.extend([f"--- {path} ---", content, ""])
    user_content = "\n".join(user_parts)

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_content},
    ]

    payload = {
        "model": _get_effective_model(os.environ.get("QWEN_MODEL", model)),
        "messages": messages,
        "temperature": 0.2,
    }

    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=180)
    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    return extract_code_from_md(content)


def call_qwen_code_complete(
    code: str,
    language_hint: str = "",
    model: str = "qwen2.5-coder-14b",
) -> str:
    """
    根据用户提供的源代码片段进行补全，返回补全后的完整代码。
    language_hint: 可选，如 "MicroPython" "C++ Arduino" 等，便于模型按语言补全。
    """
    if not (code or "").strip():
        raise ValueError("待补全的代码不能为空")
    url, headers = _get_qwen_endpoint()
    hint = f"（语言/平台：{language_hint}）" if language_hint else ""
    system = (
        "You are a code completion assistant. The user will provide a code snippet that may be "
        "incomplete. Your task is to complete it into a valid, runnable piece of code. "
        "Output ONLY the completed code, no explanations. Use a single Markdown code block "
        "(e.g. ```python or ```cpp) wrapping the entire result."
    )
    user_content = (
        f"请补全以下代码{hint}，只输出补全后的完整代码，不要解释。\n\n"
        "```\n"
        f"{code.strip()}\n"
        "```"
    )
    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": user_content},
    ]
    payload = {
        "model": _get_effective_model(os.environ.get("QWEN_MODEL", model)),
        "messages": messages,
        "temperature": 0.2,
    }
    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=120)
    content = resp.json()["choices"][0]["message"]["content"]
    return extract_code_from_md(content)


def call_qwen_code_optimize(
    code: str,
    instruction: str = "",
    model: str = "qwen2.5-coder-14b",
) -> str:
    """
    对用户提供的源代码进行优化，返回优化后的代码。
    instruction: 可选，如 "减少内存占用" "提高可读性" "适配 ESP8266" 等。
    """
    if not (code or "").strip():
        raise ValueError("待优化的代码不能为空")
    url, headers = _get_qwen_endpoint()
    extra = f"\n用户额外要求：{instruction.strip()}" if instruction.strip() else ""
    system = (
        "You are a code optimization assistant. The user will provide source code. "
        "Optimize it (performance, readability, resource usage, or as the user requests). "
        "Output ONLY the optimized code in a single Markdown code block, no explanations."
    )
    user_content = (
        "请优化以下代码，只输出优化后的完整代码，不要解释。"
        f"{extra}\n\n"
        "```\n"
        f"{code.strip()}\n"
        "```"
    )
    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": user_content},
    ]
    payload = {
        "model": _get_effective_model(os.environ.get("QWEN_MODEL", model)),
        "messages": messages,
        "temperature": 0.2,
    }
    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=120)
    content = resp.json()["choices"][0]["message"]["content"]
    return extract_code_from_md(content)


# =====================
# 电脑助手：多模式对话与终端、按路径读写文件
# =====================


def _normalize_filename_extension(name: str) -> str:
    """用户常把扩展名写成「文案docx」，补上点号变为「文案.docx」便于匹配真实文件。"""
    if not name:
        return name
    for ext in ("docx", "doc", "txt", "py", "md", "json"):
        if len(name) > len(ext) and name.lower().endswith(ext) and name[-len(ext) - 1] != ".":
            return name[:-len(ext)] + "." + ext
    return name


def resolve_file_path_from_instruction(instruction: str) -> Optional[str]:
    """
    从用户自然语言中解析出文件路径并转为绝对路径。
    支持：桌面上的xxx、桌面中某文件夹里的xxx、~/Desktop/子目录/xxx、项目下某目录里的yyy 等。
    返回绝对路径，若无法解析则返回 None。
    """
    if not (instruction or "").strip():
        return None
    text = instruction.strip()
    desktop = os.path.expanduser("~/Desktop")
    root = get_project_root()
    # 桌面(上|中) 某文件夹里(的) 文件名（如：桌面中的furina文件夹里的furina.html 或 furina.html,请输出）
    # 用 [^\s，。！？\n,;]+? 避免把句尾的半角逗号吃进文件名，导致路径不存在
    m = re.search(
        r"(?:我)?桌面(?:上|中|里)?的?\s*(\S+?)文件夹里(?:的)?\s*([^\s，。！？\n,;]+?)(?=\s*[，。！？\n,;]|$|帮|请|把|优化|修改|改|润色)",
        text,
    )
    if m:
        folder_name = m.group(1).strip().strip("""'""")
        file_name = (m.group(2) or "").strip().strip("""'""").rstrip(",;，。！？")
        if folder_name and file_name:
            file_name = _normalize_filename_extension(file_name)
            candidate = os.path.normpath(os.path.join(desktop, folder_name, file_name))
            if os.path.isfile(candidate):
                return candidate
            return candidate
    # 项目(下) 某目录(里)的 文件名
    m = re.search(
        r"项目(?:根)?(?:目录)?(?:下)?的?\s*(\S+?)(?:文件夹)?里(?:的)?\s*([^\s，。！？\n,;]+?)(?=\s*[，。！？\n,;]|$|帮|请|把|优化|修改|改|润色)",
        text,
    )
    if m:
        folder_name = m.group(1).strip().strip("""'""")
        file_name = (m.group(2) or "").strip().strip("""'""").rstrip(",;，。！？")
        if folder_name and file_name:
            file_name = _normalize_filename_extension(file_name)
            candidate = os.path.normpath(os.path.join(root, folder_name, file_name))
            if os.path.isfile(candidate):
                return candidate
            return candidate
    # 我?桌面(上|中|里)?的? xxx（如：润色桌面上的xxx.docx、桌面中的聚光日记宣传片文案docx）
    m = re.search(
        r"(?:我)?桌面(?:上|中|里)?的?\s*[：:]\s*([^\s，。！？\n]+)|"
        r"(?:我)?桌面(?:上|中|里)?的?\s*([^\s，。！？\n]+?)(?=\s*[，。！？\n]|$|帮|请|把|将|改|写|润色|修改)",
        text,
    )
    if m:
        raw_name = (m.group(1) or m.group(2) or "").strip().strip("""'""")
        if raw_name:
            name = _normalize_filename_extension(raw_name)
            candidate = os.path.normpath(os.path.join(desktop, name))
            if os.path.isfile(candidate):
                return candidate
            # 用户可能存成「文案docx」无点号，先试带点的，再试原名
            if raw_name != name:
                fallback = os.path.normpath(os.path.join(desktop, raw_name))
                if os.path.isfile(fallback):
                    return fallback
            return candidate
    # 桌面/xxx 或 桌面\xxx
    m = re.search(r"桌面[/\\]\s*([^\s，。！？\n]+)", text)
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            name = _normalize_filename_extension(name)
            return os.path.normpath(os.path.join(desktop, name))
    # ~/Desktop/xxx 或 ~/桌面/xxx（允许路径中含空格）
    m = re.search(r"~/(?:Desktop|桌面)\s*/\s*([^\s]+(?:\s+[^\s]+)*)", text)
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            name = _normalize_filename_extension(name)
            return os.path.normpath(os.path.join(desktop, name))
    m = re.search(r"~/(?:Desktop|桌面)/([^\s]+)", text)
    if m:
        return os.path.normpath(os.path.join(desktop, _normalize_filename_extension(m.group(1).strip())))
    # 项目根目录(下)的 xxx / 项目下的 xxx
    m = re.search(r"项目(?:根)?(?:目录)?(?:下)?的?\s*([^\s，。！？\n]+?)(?=\s*[，。！？\n]|$|帮|请|把|将)", text)
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            name = _normalize_filename_extension(name)
            return os.path.normpath(os.path.join(root, name))
    # 已为绝对路径或相对路径（以 . 或 / 开头，或含扩展名的单词）
    m = re.search(r"(?:^|[\s：:])((/[^\s]+)|(~/[^\s]+)|([a-zA-Z]:\\[^\s]+))", text)
    if m:
        raw = (m.group(2) or m.group(3) or m.group(4) or "").strip()
        if raw:
            return os.path.normpath(os.path.expanduser(raw))
    return None


def get_mentioned_file_paths(instruction: str) -> List[Tuple[str, str]]:
    """
    从用户指令中解析出提到的文件名（如 xxx.html、yyy.swift），在桌面与项目根下查找并返回存在的文件。
    返回 [(绝对路径, 文件名), ...]，去重且仅包含允许目录下的真实文件。
    """
    if not (instruction or "").strip():
        return []
    text = instruction.strip()
    desktop = os.path.expanduser("~/Desktop")
    root = get_project_root()
    # 匹配常见扩展名的文件名（含子路径如 furina/furina.html）
    pattern = re.compile(
        r"\b([\w.-]+/[\w.-]+\.(?:html?|swift|py|js|ts|jsx|tsx|css|txt|md|json|yaml|yml|docx|doc))\b|"
        r"\b([\w.-]+\.(?:html?|swift|py|js|ts|jsx|tsx|css|txt|md|json|yaml|yml|docx|doc))\b",
        re.IGNORECASE,
    )
    seen_names: set = set()
    candidates: List[Tuple[str, str]] = []
    for m in pattern.finditer(text):
        name = (m.group(1) or m.group(2) or "").strip().strip("'\"")
        if not name or name in seen_names:
            continue
        seen_names.add(name)
        name = _normalize_filename_extension(name)
        # 尝试多种路径
        to_try = [
            resolve_file_path_from_instruction("桌面上的 " + name),
            resolve_file_path_from_instruction("项目下的 " + name),
            os.path.normpath(os.path.join(desktop, name)),
            os.path.normpath(os.path.join(root, name)),
        ]
        # 若指令中有「某文件夹里的该文件」，主解析可能已能解析；再试桌面/项目下常见子目录
        for base in (desktop, root):
            if os.path.isdir(base):
                for entry in os.listdir(base):
                    if entry.startswith("."):
                        continue
                    sub = os.path.join(base, entry)
                    if os.path.isdir(sub):
                        to_try.append(os.path.normpath(os.path.join(sub, os.path.basename(name))))
        for path in to_try:
            if not path or not os.path.isfile(path):
                continue
            if not _is_path_under_allowed_bases(path):
                continue
            basename = os.path.basename(path)
            if not any(c[0] == path for c in candidates):
                candidates.append((path, basename))
            break
    return candidates


def _get_allowed_folder_bases() -> List[str]:
    """返回允许助手操作文件夹的根目录列表（桌面、项目根）。"""
    desktop = os.path.expanduser("~/Desktop")
    root = get_project_root()
    bases = [os.path.realpath(root)]
    if os.path.realpath(desktop) not in bases:
        bases.append(os.path.realpath(desktop))
    return bases


def _is_path_under_allowed_bases(path: str) -> bool:
    """检查 path 是否在允许的文件夹根目录之下。"""
    try:
        real = os.path.realpath(os.path.normpath(path))
    except OSError:
        return False
    for base in _get_allowed_folder_bases():
        if real == base or real.startswith(base + os.sep):
            return True
    return False


def is_path_under_allowed_bases(path: str) -> bool:
    """公开接口：检查 path 是否在允许的文件夹根目录之下（供 web 预览等使用）。"""
    return _is_path_under_allowed_bases(path)


def resolve_folder_path_from_instruction(instruction: str) -> Optional[str]:
    """
    从用户自然语言中解析出文件夹路径并转为绝对路径。
    支持：桌面上的某文件夹、项目下的某目录、~/Desktop/xxx 等。
    仅当路径存在且为目录、且在允许的根目录（桌面、项目根）之下时返回。
    """
    if not (instruction or "").strip():
        return None
    text = instruction.strip()
    desktop = os.path.expanduser("~/Desktop")
    root = get_project_root()
    # 桌面(上|中) 某名+「文件夹」→ 只取文件夹名（如：桌面上的 furina 文件夹、桌面中furina文件夹内容）
    m = re.search(
        r"(?:我)?桌面(?:上|中|里)?的?\s*(\S+?)(?:文件夹)(?=\s*|里|的|内容|$)",
        text,
    )
    if m:
        raw_name = m.group(1).strip().strip("""'""")
        if raw_name:
            candidate = os.path.normpath(os.path.join(desktop, raw_name))
            if _is_path_under_allowed_bases(candidate):
                return candidate
    # 桌面(上|里)?的? 某文件夹名（不强调扩展名）
    m = re.search(
        r"(?:我)?桌面(?:上|中|里)?的?\s*[：:]\s*([^\s，。！？\n]+)|"
        r"(?:我)?桌面(?:上|中|里)?的?\s*([^\s，。！？\n]+?)(?=\s*[，。！？\n]|$|帮|请|把|修改|编辑|批量)",
        text,
    )
    if m:
        raw_name = (m.group(1) or m.group(2) or "").strip().strip("""'""")
        if raw_name:
            candidate = os.path.normpath(os.path.join(desktop, raw_name))
            if os.path.isdir(candidate) and _is_path_under_allowed_bases(candidate):
                return candidate
            return candidate  # 允许返回路径供后续创建/列出
    m = re.search(r"桌面[/\\]\s*([^\s，。！？\n]+)", text)
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            candidate = os.path.normpath(os.path.join(desktop, name))
            if _is_path_under_allowed_bases(candidate):
                return candidate
    m = re.search(r"~/(?:Desktop|桌面)\s*/\s*([^\s]+(?:\s+[^\s]+)*)", text)
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            candidate = os.path.normpath(os.path.join(desktop, name))
            if _is_path_under_allowed_bases(candidate):
                return candidate
    m = re.search(r"~/(?:Desktop|桌面)/([^\s]+)", text)
    if m:
        candidate = os.path.normpath(os.path.join(desktop, m.group(1).strip()))
        if _is_path_under_allowed_bases(candidate):
            return candidate
    # 项目(根)?(目录)?下的? 某目录名
    m = re.search(
        r"项目(?:根)?(?:目录)?(?:下)?的?\s*([^\s，。！？\n]+?)(?=\s*[，。！？\n]|$|帮|请|把|修改|编辑|批量|里|内)",
        text,
    )
    if m:
        name = m.group(1).strip().strip("""'""")
        if name:
            candidate = os.path.normpath(os.path.join(root, name))
            if os.path.isdir(candidate) and _is_path_under_allowed_bases(candidate):
                return candidate
            # 也接受项目根本身
            if _is_path_under_allowed_bases(candidate):
                return candidate
    # 绝对路径或 ~/ 路径
    m = re.search(r"(?:^|[\s：:])((/[^\s]+)|(~/[^\s]+)|([a-zA-Z]:\\[^\s]+))", text)
    if m:
        raw = (m.group(2) or m.group(3) or m.group(4) or "").strip()
        if raw:
            candidate = os.path.normpath(os.path.expanduser(raw))
            if os.path.isdir(candidate) and _is_path_under_allowed_bases(candidate):
                return candidate
    return None


def list_directory_for_assistant(
    dir_path: str,
    pattern: str = "*",
    recursive: bool = False,
    max_entries: int = 200,
) -> Tuple[bool, List[Tuple[str, bool]], str]:
    """
    列出目录下的条目（仅允许在桌面/项目根之下）。
    返回 (ok, [(条目相对路径, 是否为文件)], error_message)。
    pattern 为 glob 风格，如 *.py；recursive 是否递归子目录。
    """
    dir_path = os.path.normpath(os.path.expanduser(dir_path.strip()))
    if not os.path.isdir(dir_path):
        return False, [], f"目录不存在：{dir_path}"
    if not _is_path_under_allowed_bases(dir_path):
        return False, [], f"仅允许操作桌面或项目根下的目录：{dir_path}"
    result: List[Tuple[str, bool]] = []
    try:
        if not recursive:
            for name in sorted(os.listdir(dir_path)):
                if name.startswith("."):
                    continue
                full = os.path.join(dir_path, name)
                rel = name
                is_file = os.path.isfile(full)
                if fnmatch.fnmatch(name, pattern):
                    result.append((rel, is_file))
                if len(result) >= max_entries:
                    break
        else:
            for root_dir, _, files in os.walk(dir_path):
                rel_base = os.path.relpath(root_dir, dir_path)
                if rel_base == ".":
                    rel_base = ""
                for f in files:
                    if f.startswith("."):
                        continue
                    rel = (os.path.join(rel_base, f) if rel_base else f).replace("\\", "/")
                    if fnmatch.fnmatch(f, pattern):
                        result.append((rel, True))
                    if len(result) >= max_entries:
                        return True, result, ""
    except OSError as e:
        return False, [], str(e)
    return True, result, ""


def read_folder_files_for_assistant(
    dir_path: str,
    pattern: str = "*",
    max_files: int = 20,
    recursive: bool = False,
) -> Tuple[bool, List[Tuple[str, str]], str]:
    """
    读取目录下匹配 pattern 的文件内容（仅允许在桌面/项目根之下）。
    返回 (ok, [(相对路径, 文件内容)], error_message)。
    """
    dir_path = os.path.normpath(os.path.expanduser(dir_path.strip()))
    if not _is_path_under_allowed_bases(dir_path):
        return False, [], f"仅允许操作桌面或项目根下的目录：{dir_path}"
    ok, entries, err = list_directory_for_assistant(dir_path, pattern=pattern, recursive=recursive, max_entries=max_files * 2)
    if not ok:
        return False, [], err
    file_entries = [(rel, is_file) for rel, is_file in entries if is_file][:max_files]
    result: List[Tuple[str, str]] = []
    for rel, _ in file_entries:
        full = os.path.join(dir_path, rel)
        if not os.path.isfile(full):
            continue
        ok_read, content, err_read = read_file_content_for_assistant(full)
        if not ok_read:
            return False, [], f"读取 {rel} 失败：{err_read}"
        result.append((rel, content))
    return True, result, ""


def write_assistant_results_to_folder(
    dir_path: str,
    file_edits: dict,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> Tuple[bool, List[str]]:
    """
    将助手返回的多文件内容写回目录（仅允许在桌面/项目根之下）。
    file_edits 为 { "相对路径": "内容", ... }，路径统一用 /。
    progress_callback(msg) 可选，用于汇报行动进程（如「开始创建 xxx」「xxx 创建完成」）。
    返回 (全部成功, [错误信息列表])。
    """
    dir_path = os.path.normpath(os.path.expanduser(dir_path.strip()))
    if not _is_path_under_allowed_bases(dir_path):
        return False, [f"目录不允许写入：{dir_path}"]
    errors: List[str] = []
    for rel, content in file_edits.items():
        rel = rel.replace("\\", "/").lstrip("/")
        if ".." in rel or rel.startswith("/"):
            errors.append(f"非法路径：{rel}")
            continue
        full = os.path.join(dir_path, rel)
        try:
            real_full = os.path.realpath(full)
            real_dir = os.path.realpath(dir_path)
            if real_full != real_dir and not real_full.startswith(real_dir + os.sep):
                errors.append(f"路径越界：{rel}")
                continue
        except OSError:
            errors.append(f"路径无效：{rel}")
            continue
        name = os.path.basename(rel)
        if progress_callback and name:
            progress_callback("开始创建 " + name)
        ok, err = write_assistant_result_to_file(full, content)
        if progress_callback and name:
            progress_callback(name + " 创建完成")
        if not ok:
            errors.append(f"{rel}: {err}")
    return len(errors) == 0, errors


def resolve_create_target_from_instruction(instruction: str) -> Optional[Tuple[str, str]]:
    """
    从用户指令中解析「创建」目标：在桌面或项目根下创建文件夹。
    返回 (基础目录绝对路径, 文件夹名) 或 None。仅允许桌面、项目根。
    """
    if not (instruction or "").strip():
        return None
    # 规范化：全角空格、多余空白
    text = (instruction or "").strip().replace("\u3000", " ")
    text = re.sub(r"\s+", " ", text)
    desktop = os.path.expanduser("~/Desktop")
    root = get_project_root()
    folder_name = None
    base = desktop
    # 【创造 Agent】在桌面创建名为「XXX」→ 优先取引号内名称，避免被「创建 xxx 的 文件夹」误解析为「名为开发者」
    m = re.search(r"(?:在\s*桌面\s*)?创建\s*(?:名为)?\s*[「『]\s*([^」』\s，。！？]+)\s*[」』]", text)
    if m:
        raw = (m.group(1) or "").strip()
        raw = re.sub(r"\s+", "_", re.sub(r"[^\w\u4e00-\u9fff\-_\s]", "", raw))
        if raw:
            folder_name = (raw or "新创建")[:50]
    # 做一个 xxx 的网站 / 帮我做一个 furina 的网站 / 原神重云的介绍网站（「的」与「网站」之间可有「介绍」等词）
    if not folder_name:
        m = re.search(r"(?:帮我)?做(?:一个)?\s*[「『]?\s*([^」』\s，。！？]+?)\s*[」』]?\s*的.*?网站", text)
        if m:
            raw = re.sub(r"\s+", "_", (m.group(1) or "").strip())
            raw = re.sub(r"[^\w\u4e00-\u9fff\-_]", "", raw)
            folder_name = (raw or "网站") + "_网站"
    # 创建 xxx 的网站/文件夹 / 在桌面创建 xxx
    if not folder_name:
        m = re.search(r"创建\s*[「『]?\s*([^」』\s，。！？]+)\s*[」』]?\s*的\s*(?:网站|文件夹)|在\s*桌面\s*创建\s*[「『]?\s*([^」』\s，。！？]+)", text)
        if m:
            raw = (m.group(1) or m.group(2) or "").strip()
            raw = re.sub(r"\s+", "_", re.sub(r"[^\w\u4e00-\u9fff\-_\s]", "", raw))
            folder_name = (raw or "新创建")[:50]
    # 做 xxx 网站（无「的」）
    if not folder_name:
        m = re.search(r"(?:帮我)?做(?:一个)?\s*([^\s，。！？]+)\s*网站", text)
        if m:
            raw = re.sub(r"\s+", "_", re.sub(r"[^\w\u4e00-\u9fff\-_\s]", "", (m.group(1) or "").strip()))
            folder_name = (raw or "网站") + "_网站"
    # 做 xxx 的? 网页 / 做 xxx 网页小游戏 / 做 xxx 小游戏（如：帮我做一个抢红包的网页小游戏）
    if not folder_name and re.search(r"做.*(?:网页|小游戏)", text):
        m = re.search(r"(?:帮我)?做(?:一个)?\s*[「『]?\s*([^」』\s，。！？]+?)\s*[」』]?\s*的?\s*(?:网页|小游戏)", text)
        if m:
            raw = (m.group(1) or "").strip().rstrip("的").strip()
            raw = re.sub(r"\s+", "_", raw)
            raw = re.sub(r"[^\w\u4e00-\u9fff\-_]", "", raw)
            if raw:
                folder_name = (raw or "网页") + "_网页"
    # 做/开发 xxx 的 (iOS/Android)? 软件/应用/App（如：帮我做一个写日记的ios软件、写日记的ios 软件、开发一个记账App）
    if not folder_name:
        m = re.search(
            r"(?:帮我)?(?:做|开发)(?:一个)?\s*[「『]?\s*([^」』\s，。！？]+?)\s*[」』]?\s*的?\s*(?:ios\s*|android\s*|安卓\s*)?(?:软件|应用|App)",
            text,
            re.IGNORECASE,
        )
        if m:
            raw = re.sub(r"\s+", "_", (m.group(1) or "").strip())
            raw = re.sub(r"[^\w\u4e00-\u9fff\-_]", "", raw)
            folder_name = (raw or "应用") + "_项目"
    # 更宽松：只要包含「做」+「软件/应用/App」且能截出主题名（如「写日记的ios 软件」）
    if not folder_name and re.search(r"做.*(?:软件|应用|App)", text, re.IGNORECASE):
        m = re.search(r"(?:帮我)?做(?:一个)?\s*(.+?)(?:的)?\s*(?:ios|android|安卓)?\s*(?:软件|应用|App)", text, re.IGNORECASE)
        if m:
            raw = (m.group(1) or "").strip().rstrip("的").strip()
            raw = re.sub(r"\s+", "_", raw)
            raw = re.sub(r"[^\w\u4e00-\u9fff\-_]", "", raw)
            if raw:
                folder_name = (raw or "应用") + "_项目"
    if not folder_name:
        return None
    folder_name = folder_name.strip("_") or "新创建"
    if not _is_path_under_allowed_bases(desktop):
        return (root, folder_name)
    return (desktop, folder_name)


def ensure_directory_and_write_files(
    parent_dir: str,
    file_edits: dict,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> Tuple[bool, List[str]]:
    """
    确保父目录存在（可多级创建），然后写入多文件。仅允许在桌面/项目根之下。
    progress_callback(msg) 可选，用于汇报行动进程。
    返回 (全部成功, [错误信息列表])。
    """
    parent_dir = os.path.normpath(os.path.expanduser(parent_dir.strip()))
    if not _is_path_under_allowed_bases(parent_dir):
        return False, [f"仅允许在桌面或项目根下创建：{parent_dir}"]
    folder_name = os.path.basename(parent_dir.rstrip(os.sep)) or "项目"
    try:
        if progress_callback:
            progress_callback("开始创建 " + folder_name + " 文件夹")
        os.makedirs(parent_dir, exist_ok=True)
        if progress_callback:
            progress_callback(folder_name + " 文件夹创建完成")
        for rel in file_edits:
            subdir = os.path.dirname(os.path.join(parent_dir, rel.replace("\\", "/")))
            if subdir and subdir != parent_dir:
                subdir_name = os.path.basename(subdir.rstrip(os.sep))
                if progress_callback and subdir_name:
                    progress_callback("开始创建 " + subdir_name + " 文件夹")
                os.makedirs(subdir, exist_ok=True)
                if progress_callback and subdir_name:
                    progress_callback(subdir_name + " 文件夹创建完成")
    except OSError as e:
        return False, [str(e)]
    return write_assistant_results_to_folder(parent_dir, file_edits, progress_callback=progress_callback)


# 电脑助手：所有文件读写均通过终端命令执行
def _run_command_for_assistant(
    cmd: List[str],
    stdin: Optional[str] = None,
    cwd: Optional[str] = None,
    timeout: int = 30,
) -> Tuple[bool, str, str]:
    """执行一条命令（用于助手读/写文件等），返回 (ok, stdout, stderr)。"""
    try:
        proc = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout,
            cwd=cwd or get_project_root(),
            input=stdin,
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        return proc.returncode == 0, out, err
    except subprocess.TimeoutExpired:
        return False, "", f"命令执行超时（{timeout}s）"
    except Exception as e:
        return False, "", str(e)


# .docx 读写（仅用标准库 zipfile + xml，无需 python-docx）
_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _read_docx_stdlib(path: str) -> str:
    """仅用标准库从 .docx 中提取正文文本（按段落用换行连接）。"""
    with zipfile.ZipFile(path, "r") as z:
        data = z.read("word/document.xml")
    root = ET.fromstring(data)
    paragraphs = []
    for p in root.findall(f".//{{{_W_NS}}}p"):
        texts = []
        for t in p.findall(f".//{{{_W_NS}}}t"):
            if t.text:
                texts.append(t.text)
            if t.tail:
                texts.append(t.tail)
        paragraphs.append("".join(texts))
    return "\n".join(paragraphs)


def _write_docx_stdlib(path: str, content: str) -> None:
    """仅用标准库将纯文本写入为 .docx（多段按 \\n\\n 分割，无 python-docx）。"""
    def escape(s: str) -> str:
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    blocks = [b.strip() for b in (content or "").split("\n\n") if b.strip()]
    if not blocks:
        blocks = [""]
    para_xml = "".join(
        f'<w:p><w:r><w:t xml:space="preserve">{escape(bl)}</w:t></w:r></w:p>'
        for bl in blocks
    )
    doc_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:document xmlns:w="{_W_NS}">'
        "<w:body>"
        f"{para_xml}"
        '<w:sectPr><w:pgSz w:w="12240" w:h="15840"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/><w:docGrid w:linePitch="360"/></w:sectPr>'
        "</w:body></w:document>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    doc_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/_rels/document.xml.rels", doc_rels)
        z.writestr("word/document.xml", doc_xml)


# 供终端执行的 .docx 读/写脚本内容（内联，通过 temp 文件运行）
_DOCX_READ_SCRIPT = r"""
import sys, zipfile, xml.etree.ElementTree as ET
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
path = sys.argv[1]
with zipfile.ZipFile(path, "r") as z:
    data = z.read("word/document.xml")
root = ET.fromstring(data)
for p in root.findall(".//{%s}p" % W_NS):
    for t in p.findall(".//{%s}t" % W_NS):
        if t.text: print(t.text, end="")
        if t.tail: print(t.tail, end="")
    print()
"""

_DOCX_WRITE_SCRIPT = r"""
import sys, zipfile, re
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
path = sys.argv[1]
content = sys.stdin.read()
def escape(s): return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
blocks = [b.strip() for b in content.split("\n\n") if b.strip()] or [""]
para_xml = "".join('<w:p><w:r><w:t xml:space="preserve">%s</w:t></w:r></w:p>' % escape(bl) for bl in blocks)
doc_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><w:document xmlns:w="%s"><w:body>%s<w:sectPr><w:pgSz w:w="12240" w:h="15840"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/><w:docGrid w:linePitch="360"/></w:sectPr></w:body></w:document>' % (W_NS, para_xml)
ct = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>'
rels = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>'
doc_rels = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
    z.writestr("[Content_Types].xml", ct)
    z.writestr("_rels/.rels", rels)
    z.writestr("word/_rels/document.xml.rels", doc_rels)
    z.writestr("word/document.xml", doc_xml)
"""

# 供终端执行的 .pptx 写入脚本（依赖 python-pptx：pip install python-pptx）
# 约定：stdin 为纯文本，按双换行分页，每页首行为标题，其余行为正文（以 - 或 * 开头视为 bullet）
# 图片：单独一行 [IMG: url] 或 [图片: url] 或 [IMG: /本地路径] 会下载/读取并插入到当前页（可多行多图）
_PPTX_WRITE_SCRIPT = r"""
import sys
import re
import os
import tempfile
import urllib.request
path = sys.argv[1]
content = sys.stdin.read()
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    sys.stderr.write("请先安装: pip install python-pptx\n")
    sys.exit(1)
IMG_PATTERN = re.compile(r'^\s*\[(?:IMG|图片)\s*:\s*(.+?)\]\s*$', re.IGNORECASE)
def is_img_line(line):
    m = IMG_PATTERN.match(line.strip())
    return m.group(1).strip() if m else None
def get_image_path(ref):
    ref = ref.strip()
    if not ref:
        return None
    if ref.startswith(('http://', 'https://')):
        try:
            req = urllib.request.Request(ref, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=15) as resp:
                ext = '.png' if 'png' in (resp.headers.get('Content-Type') or '') else '.jpg'
                fd, loc = tempfile.mkstemp(suffix=ext)
                with os.fdopen(fd, 'wb') as f:
                    f.write(resp.read())
                return loc
        except Exception:
            return None
    if os.path.isfile(ref):
        return ref
    return None
raw_slides = [s.strip() for s in (content or "").split("\n\n") if s.strip()]
if not raw_slides:
    raw_slides = [""]
prs = Presentation()
title_layout = prs.slide_layouts[0]
content_layout = prs.slide_layouts[1]
for i, block in enumerate(raw_slides):
    lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
    title_text = lines[0] if lines else ""
    body_lines = lines[1:] if len(lines) > 1 else []
    img_refs = []
    text_lines = []
    for ln in body_lines:
        ref = is_img_line(ln)
        if ref:
            img_refs.append(ref)
        else:
            text_lines.append(ln)
    if i == 0 and not text_lines and not img_refs:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title_text or "标题"
    else:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title_text or "标题"
        tf = slide.placeholders[1].text_frame
        if text_lines:
            tf.text = (text_lines[0].lstrip("-* ").strip() or text_lines[0])[:4096]
            for line in text_lines[1:]:
                p = tf.add_paragraph()
                p.text = (line.lstrip("-* ").strip() or line)[:4096]
                p.level = 1 if line.strip().startswith(("-", "*")) else 0
        else:
            tf.text = ""
        for idx, ref in enumerate(img_refs[:4]):
            img_path = get_image_path(ref)
            if not img_path:
                continue
            try:
                left = Inches(6) + (idx % 2) * Inches(3.2)
                top = Inches(1.8) + (idx // 2) * Inches(2.4)
                slide.shapes.add_picture(img_path, left, top, width=Inches(3))
            except Exception:
                pass
            if ref.startswith(('http://', 'https://')):
                try:
                    os.unlink(img_path)
                except Exception:
                    pass
prs.save(path)
"""


def read_file_content_for_assistant(path: str) -> Tuple[bool, str, str]:
    """
    通过终端命令读取文件内容供助手使用。支持文本及 .docx。
    返回 (ok, content, error_message)。
    """
    path = os.path.normpath(os.path.expanduser(path.strip()))
    if not os.path.isfile(path):
        return False, "", f"文件不存在：{path}"
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        fd, script_path = tempfile.mkstemp(suffix=".py", prefix="lumi_read_docx_")
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as f:
                f.write(_DOCX_READ_SCRIPT)
            ok, out, err = _run_command_for_assistant(
                [sys.executable, script_path, path],
                timeout=15,
            )
            if not ok and "BadZipFile" in err:
                return False, "", "该文件不是有效的 .docx（ZIP 格式）"
            return ok, out, err
        finally:
            try:
                os.unlink(script_path)
            except OSError:
                pass
    ok, out, err = _run_command_for_assistant(
        [sys.executable, "-c", "import sys; print(open(sys.argv[1], encoding='utf-8', errors='replace').read())", path],
        timeout=10,
    )
    return ok, out, err


def extract_content_to_write_from_reply(reply: str) -> str:
    """从助手回复中提取将要写回文件的内容（与 write 时使用的逻辑一致）。"""
    content = (reply or "").strip()
    code_blocks = re.findall(r"```(?:\w*)\s*([\s\S]*?)```", content)
    if code_blocks:
        content = max((b.strip() for b in code_blocks), key=len)
    return content


def extract_html_from_reply(reply: str) -> Optional[str]:
    """当回复中无 ---FILE:--- 时，若包含整段 HTML（<!DOCTYPE 或 <html ... </html>），则截取返回。"""
    if not (reply or "").strip():
        return None
    text = (reply or "").strip()
    start = text.find("<!DOCTYPE")
    if start < 0:
        start = text.find("<html")
    if start < 0:
        return None
    end = text.rfind("</html>")
    if end < 0 or end <= start:
        return None
    return text[start : end + 7].strip()


def extract_run_command_from_reply(reply: str) -> Optional[str]:
    """
    从助手回复中解析需由系统执行的终端命令（仅第一条）。
    约定格式：---RUN: 单行命令 ---（模型根据用户输入判断是否需要执行终端时输出此格式）。
    若无该格式但整条回复为单行且以白名单命令开头，也视为要执行的命令（兼容模型只输出命令文本的情况）。
    返回命令字符串，若无则返回 None。
    """
    commands = extract_run_commands_from_reply(reply)
    return commands[0] if commands else None


def extract_run_commands_from_reply(reply: str) -> List[str]:
    """
    从助手回复中解析所有 ---RUN: 命令 ---，按出现顺序返回列表。
    复杂项目时 AI 可输出多条 ---RUN:---，系统会依次在项目目录下执行（创建文件夹、安装依赖、编译等）。
    """
    if not (reply or "").strip():
        return []
    text = (reply or "").strip()
    pattern = re.compile(r"---RUN:\s*([^\n\-]+?)\s*---", re.IGNORECASE)
    commands = []
    for m in pattern.finditer(text):
        cmd = m.group(1).strip()
        if cmd:
            commands.append(cmd)
    if commands:
        return commands
    # 兼容：整条回复只有一行且以常见“只读/查信息”命令开头时，视为要执行的命令
    if "\n" not in text or text.count("\n") == 0:
        first_word = (text.split() or [""])[0].lower()
        if first_word in (
            "system_profiler", "uname", "sw_vers", "whoami", "hostname", "date",
            "ls", "pwd", "cat", "echo", "df", "top", "ps", "uptime", "env", "printenv",
            "lscpu", "vm_stat", "sysctl", "free", "ifconfig", "ip", "netstat", "nslookup", "dig",
            "head", "tail", "wc", "file", "stat", "du", "which", "whereis", "find",
            "grep", "rg", "diff", "tree", "ping", "curl", "wget",
        ):
            return [text]
    return []


def write_assistant_result_to_file(
    path: str,
    content: str,
    original_length: Optional[int] = None,
) -> Tuple[bool, str]:
    """
    通过终端命令将助手返回的内容写回文件。支持文本及 .docx。
    original_length: 若传入原文件长度，当待写内容过短（< 原长度 50%）时拒绝写入，防止用片段覆盖整文件。
    返回 (ok, error_message)。
    """
    path = os.path.normpath(os.path.expanduser(path.strip()))
    content = extract_content_to_write_from_reply(content)
    if original_length is not None and original_length > 0:
        if len(content) < original_length * 0.5:
            return (
                False,
                "返回内容过短，可能仅为修改片段，已拒绝写入以保护原文件。请要求助手「输出完整文件内容」后重试。",
            )
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        fd, script_path = tempfile.mkstemp(suffix=".py", prefix="lumi_write_docx_")
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as f:
                f.write(_DOCX_WRITE_SCRIPT)
            ok, out, err = _run_command_for_assistant(
                [sys.executable, script_path, path],
                stdin=content,
                timeout=15,
            )
            return ok, err
        finally:
            try:
                os.unlink(script_path)
            except OSError:
                pass
    if ext == ".pptx":
        fd, script_path = tempfile.mkstemp(suffix=".py", prefix="lumi_write_pptx_")
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as f:
                f.write(_PPTX_WRITE_SCRIPT)
            ok, out, err = _run_command_for_assistant(
                [sys.executable, script_path, path],
                stdin=content,
                timeout=30,
            )
            return ok, err
        finally:
            try:
                os.unlink(script_path)
            except OSError:
                pass
    # Windows 上直接使用 Python 文件操作，避免命令行参数问题
    try:
        # 确保目录存在
        dir_path = os.path.dirname(path)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)
        # 直接写入文件
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return True, ""
    except Exception as e:
        # 如果直接写入失败，回退到命令方式
        ok, out, err = _run_command_for_assistant(
            [sys.executable, "-c", "import sys; open(sys.argv[1],'w',encoding='utf-8').write(sys.stdin.read())", path],
            stdin=content,
            timeout=10,
        )
        return ok, err or str(e)


ASSISTANT_MODES = (
    "polish",           # 文本润色
    "edit_code",        # 修改代码
    "complete_code",    # 补全代码
    "terminal",         # 终端命令（生成或执行）
    "folder_edit",      # 文件夹内多文件批量修改
    "list_folder",      # 查看文件夹内容
    "create_file",      # 自主创建文件/文件夹（如做网站）
    "todo",             # TODO 规划
    "plan",             # Plan 模式
    "deep_think",      # 深度思考
    "custom_command",   # 自定义 Command
)


def infer_assistant_mode(instruction: str, context: Optional[dict] = None) -> str:
    """
    根据用户输入的文本（及可选的文件路径等上下文）自动推断应使用的助手模式。
    返回 ASSISTANT_MODES 之一，默认 deep_think。
    """
    if not (instruction or "").strip():
        return "deep_think"
    text = (instruction or "").strip().lower()
    ctx = context or {}
    file_path = (ctx.get("file_path") or "").strip()
    ext = os.path.splitext(file_path)[1].lower() if file_path else ""

    # 自主创建文件/网站/软件/网页/小游戏（优先于 list_folder）
    if re.search(
        r"(?:帮我)?做(?:一个)?\s*.+\s*的?\s*网站|创建\s*.+\s*的\s*(?:网站|文件夹)|在\s*桌面\s*创建|"
        r"(?:帮我)?(?:做|开发)(?:一个)?\s*.+\s*的?\s*(?:ios\s*|android\s*|安卓\s*)?(?:软件|应用|App)|"
        r"(?:帮我)?做(?:一个)?\s*.+(?:网页|小游戏)",
        text,
        re.IGNORECASE,
    ):
        return "create_file"
    # 查看/列出文件夹内容
    if re.search(
        r"查看|列出|里有什么|有什么文件|(?:文件夹|目录)(?:里)?的?内容",
        text,
    ):
        return "list_folder"
    # 终端命令：运行命令/脚本、终端、pip、ls、cd 等
    if re.search(
        r"运行(?:命令|脚本)?|执行(?:命令|终端)|终端|(?:写个|给出)(?:命令|脚本)|cmd|pip\s|pip3\s|^ls\s|列出(?:文件)?|cd\s|跑一下",
        text,
    ):
        return "terminal"
    # TODO / 待办
    if re.search(r"todo|待办|任务列表|列一下要做|记一下要做|代办|待办事项", text):
        return "todo"
    # 计划 / 步骤
    if re.search(r"计划|规划|步骤|分步|怎么做|如何实现|第一步|plan\s|分几步", text):
        return "plan"
    # 自定义 Command
    if re.search(r"自定义\s*command|执行\s*.\s*命令", text):
        return "custom_command"
    # 补全代码
    if re.search(r"补全|补全代码|继续写|写下去|complete\s*code|接着写", text):
        return "complete_code"
    # 文件夹内批量修改：明确「所有」「批量」时用 folder_edit；「某文件夹里的某文件」已在上面解析为单文件，不会进这里
    if re.search(
        r"批量(?:修改|编辑)|(?:里|下)的?所有(?:\.py|\.txt|文件)?|整个目录(?:下)?(?:的)?(?:所有)?|目录里(?:的)?所有",
        text,
    ):
        return "folder_edit"
    if re.search(r"(?:文件夹|目录)(?:下|里|内)(?:的)?所有|(?:文件夹|目录)下(?:的)?(?:所有)?文件", text):
        return "folder_edit"
    # 修改代码：结合文件扩展名或“代码”关键词
    if re.search(
        r"修改代码|改代码|编辑代码|refactor|fix\s*bug|修\s*bug|在\s*.\s*里加|改写代码|改\s*main\.py|修改\s*.*\.py|\.js\s|\.ts\s|\.html\s",
        text,
    ):
        return "edit_code"
    if ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css", ".cpp", ".c", ".h") and re.search(
        r"修改|改|编辑|优化|fix|refactor|润色代码", text
    ):
        return "edit_code"
    # 文本润色：文案、docx、润色 等；或已解析到 .docx 文件
    if ext in (".docx", ".doc", ".txt") and re.search(r"润色|改写|修改|改", text):
        return "polish"
    if re.search(
        r"润色|改写(?:文案)?|改文案|修改文案|宣传片文案|文档润色|写得更流畅|改得更通顺|\.docx|文案\s*docx",
        text,
    ):
        return "polish"
    # 深度思考：分析、为什么、推理
    if re.search(r"深度思考|分析一下|想一想|为什么|原因是什么|推理|解释一下", text):
        return "deep_think"
    # 默认：通用问题用深度思考
    return "deep_think"


def call_qwen_assistant(
    mode: str,
    instruction: str,
    context: Optional[dict] = None,
    model: str = "qwen2.5-coder-14b",
) -> str:
    """
    电脑助手统一入口：按 mode 使用不同 system prompt，返回模型回复。
    context 可选：code, file_path, selected_text, custom_command 等。
    """
    url, headers = _get_qwen_endpoint()
    context = context or {}

    # 当带文件内容时，说明系统已读取本地文件，模型只需输出修改后内容
    file_edit_note = ""
    if context.get("file_content") is not None:
        file_edit_note = (
            " The user's local file content has been read by the system and attached below. "
            "You ONLY output the modified/polished full content—no explanation, no summary before or after. "
            "The system will reply to the user with one short sentence (e.g. 已经帮你改好了). "
            "Do NOT say you cannot access files—just output the result."
        )
    system_prompts = {
        "polish": (
            "You are a professional copy editor. The user will give you text to polish. "
            "Output ONLY the polished text in the same language, more fluent and clear. No explanations."
            + file_edit_note
        ),
        "edit_code": (
            "You are a senior programmer. The user will describe how to modify code and may provide code. "
            "Output the complete modified code in a single Markdown code block. If no code is given, output a complete script. "
            "Use the same language as the input. No extra explanation outside the block. "
            "CRITICAL: You MUST output the ENTIRE file content. Never output only a snippet or the changed part—the system will replace the whole file with your output; missing parts will be lost."
            + file_edit_note
        ),
        "complete_code": (
            "You are a code completion expert. The user will provide incomplete code. "
            "Output ONLY the completed code in a single Markdown code block. Same language and style."
        ),
        "terminal": (
            "You are a CLI assistant. The user will describe a task to do on their computer (e.g. list files, run Python, compress). "
            "Output ONLY the exact terminal command(s) to run, one per line if multiple. "
            "Use safe, common commands. Prefer Python scripts for complex tasks. No explanations."
        ),
        "folder_edit": (
            "You are a senior programmer. The user wants to modify MULTIPLE files in a folder. "
            "You are given the folder path and the content of each file. "
            "Output the modified content for each file using EXACTLY this format (no other text or explanation):\n"
            "---FILE: relative/path/filename---\n"
            "(full file content)\n"
            "---FILE: next_filename---\n"
            "(content)\n"
            "Rules: 1) Use only ---FILE: path--- blocks; 2) Path must be the same relative path as given (e.g. a.py, src/b.py); "
            "3) Output ALL files that were provided, each exactly once; 4) Apply the user's requested change to each file."
        ),
        "list_folder": (
            "You are a helpful assistant. The user asked to view the contents of a folder. "
            "You will be given the folder path and a listing of its contents (files and subfolders). "
            "Reply in Chinese with a clear, concise list or summary of what is in the folder. No need to output commands."
        ),
        "create_file": (
            "You are a senior developer. The user wants you to CREATE a new project (website, 软件/App, or Python Agent). "
            "CRITICAL: Do NOT output any design plan (设计方案), 思路, 步骤说明, or text like '让我一步步思考'. "
            "Your reply must start with ---FILE: path--- and contain ONLY ---FILE:--- blocks and optionally ---RUN:--- lines. "
            "1) Output all project files using EXACTLY this format:\n"
            "---FILE: relative/path/filename---\n"
            "(full file content)\n"
            "---FILE: next_path---\n"
            "(content)\n"
            "The system will create any folders needed (e.g. src/main.py creates src/). Paths are relative to the project folder. "
            "2) After all ---FILE:--- blocks, you MAY output one or more terminal commands for setup/build/install. Each on its own line:\n"
            "---RUN: exact_shell_command ---\n"
            "These run in order, in the created project directory. Examples: ---RUN: pip install -r requirements.txt ---, ---RUN: npm install ---, ---RUN: pio run ---, ---RUN: mkdir -p build ---. "
            "3) If the user asks for a 网站: output HTML/CSS/JS, must include index.html. "
            "4) If the user asks for 软件/App: output the project in ---FILE: path--- blocks; iOS (Swift/SwiftUI), Android (Kotlin), or cross-platform (React Native/Flutter). "
            "5) If the instruction starts with 【创造 Agent】 or is about creating an Agent: you MUST output a Python-based agent project WITH a graphical user interface (GUI). Include: main.py (or run.py) as entry point, agent logic (e.g. LLM/client calls, tool use, loop), a GUI module (use tkinter, PyQt/PySide, Gradio, Streamlit, or Flask/FastAPI with HTML/CSS/JS for web UI), requirements.txt with needed packages (e.g. openai, requests, and GUI libraries), config.py with API provider auto-detection, optional .env.example, README.md. The GUI must allow users to interact with the agent (e.g. input, output display, controls). CRITICAL: In config.py, implement automatic API provider detection based on the API key or environment variables. Support: OpenAI (base_url='https://api.openai.com/v1', models like 'gpt-3.5-turbo'), DeepSeek (base_url='https://api.deepseek.com', models like 'deepseek-chat'), DashScope/Aliyun (base_url='https://dashscope.aliyuncs.com/compatible-mode/v1', models like 'qwen-turbo'), Moonshot (base_url='https://api.moonshot.cn/v1', models like 'moonshot-v1-8k'). Detect provider from API key prefix, environment variables (DEEPSEEK_API_KEY, DASHSCOPE_API_KEY, MOONSHOT_API_KEY), or API_PROVIDER env var, then set the correct base_url and default model. Do NOT output a website or front-end only project for Agent requests. "
            "6) First line of your response must be ---FILE: or the first file path; no preamble. Do NOT show code in chat (no ```); deliver only via ---FILE:--- and ---RUN:---."
        ),
        "todo": (
            "You are a task planner. The user will describe a goal. "
            "Output a concise TODO list in this format, one per line: \"- [ ] item\" or \"- [x] item\" for done. "
            "Use Chinese. No other text."
        ),
        "plan": (
            "You are a step-by-step planner. The user will describe a task. "
            "Output a clear numbered plan (1. 2. 3. ...) in Chinese. Be concise and actionable."
        ),
        "deep_think": (
            "You are a careful reasoner and assistant. Based on the user's intent, you may:\n"
            "1) Reply with text (analysis, plan, explanation) when they ask for discussion or ideas.\n"
            "2) Create files by outputting ONLY ---FILE: relative/path--- blocks when they clearly ask you to create/make something (e.g. 做xxx网站/软件/网页/小游戏). Use format:\n"
            "---FILE: path---\n(content)\n---FILE: next---\n(content)\n"
            "The system will create any folders needed. After ---FILE:--- blocks you may add one or more ---RUN: exact_shell_command --- (one per line); they run in order in the created project directory (e.g. pip install -r requirements.txt, npm install, pio run).\n"
            "For PPT/幻灯片: use ---FILE: name.pptx --- with plain text; [IMG: url] for images. No design plan or preamble when creating—output ---FILE: blocks directly.\n"
            "3) When the user asks to run a terminal command or view system info (e.g. 通过终端查看电脑信息、执行命令), output ---RUN: exact_shell_command ---. You may output multiple ---RUN:--- lines; the system runs each and shows output. Do NOT say you cannot run the terminal.\n"
            "Use Chinese for normal replies. CRITICAL: Do NOT show code in chat (no ```). Deliver code only via ---FILE:--- and ---RUN:---; the system will save files and run commands."
        ),
        "custom_command": (
            "You are an assistant that executes user-defined commands. The user will give a command name and/or description. "
            "Interpret it and output: either a concrete action (script, command, or text result) in Chinese. "
            "If it's a script, output the full script in a Markdown code block."
        ),
    }
    system = system_prompts.get(mode) or system_prompts["deep_think"]

    user_parts = [instruction]
    if mode == "create_file":
        user_parts.append(
            "\n【必须】请直接以 ---FILE: 相对路径--- 开头输出多文件完整代码，不要输出任何设计方案、思路、步骤说明或「让我一步步思考」等文字。"
        )
        if (instruction or "").strip().startswith("【创造 Agent】") or (instruction or "").strip().startswith("【创造Agent】"):
            user_parts.append(
                "\n【创造 Agent 必须】本项目必须是基于 Python 的 Agent 项目，且必须包含图形化操作界面（GUI）。须包含：入口脚本（如 main.py 或 run.py）、Agent 核心逻辑（如调用大模型/API、工具调用、对话循环）、图形化界面模块（可使用 tkinter、PyQt/PySide、Gradio、Streamlit，或 Flask/FastAPI + HTML/CSS/JS 构建 Web 界面）、requirements.txt（需包含 GUI 相关依赖）、config.py（必须实现 API 提供商自动识别功能）、可选 .env.example、README.md。界面需支持用户与 Agent 交互（如输入框、输出显示、控制按钮等）。\n\n【API 提供商自动识别要求】在 config.py 中必须实现自动识别 API 提供商的功能：\n1. 支持识别：OpenAI（base_url='https://api.openai.com/v1'，默认模型如 'gpt-3.5-turbo'）、DeepSeek（base_url='https://api.deepseek.com'，默认模型如 'deepseek-chat'）、阿里云 DashScope（base_url='https://dashscope.aliyuncs.com/compatible-mode/v1'，默认模型如 'qwen-turbo'）、Moonshot（base_url='https://api.moonshot.cn/v1'，默认模型如 'moonshot-v1-8k'）。\n2. 识别方式：优先检查环境变量（DEEPSEEK_API_KEY、DASHSCOPE_API_KEY、MOONSHOT_API_KEY、OPENAI_API_KEY）或 API_PROVIDER 环境变量；若未设置，可根据 API key 特征或默认使用 DeepSeek。\n3. 根据识别结果自动设置正确的 base_url 和默认模型名称。\n4. 在 agent.py 中使用 OpenAI 客户端时，必须传入 base_url 参数（如 client = OpenAI(api_key=api_key, base_url=detected_base_url)）。\n不要生成纯命令行或无界面的 Agent 项目。"
            )
    if context.get("file_content") is not None:
        user_parts.append(
            "\n【重要】系统已把用户电脑上的该文件内容附在下方，你已能直接看到，请基于该内容修改并直接输出完整结果。"
            "切勿说「无法直接访问」「无法查看你的文件」等，只需输出修改后的完整文件内容。\n"
            "当前文件内容：\n```\n" + (context["file_content"] or "").strip() + "\n```"
        )
    folder_files = context.get("folder_files")  # list of (rel_path, content)
    if folder_files:
        user_parts.append("\n【说明】系统已读取文件夹内以下文件，请按用户要求修改每个文件，并用 ---FILE: 相对路径--- 格式输出全部文件。\n")
        user_parts.append("文件夹路径：" + str(context.get("folder_path", "")))
        for rel, content in folder_files:
            user_parts.append("\n--- 文件：" + rel + " ---\n" + (content or "").strip())
    folder_listing = context.get("folder_listing")  # str 或 list，供查看文件夹内容
    if folder_listing is not None:
        if isinstance(folder_listing, (list, tuple)):
            lines = []
            for item in folder_listing:
                if isinstance(item, (list, tuple)):
                    name, is_file = item[0], item[1] if len(item) > 1 else True
                    lines.append("  " + name + (" (文件)" if is_file else " (文件夹)"))
                else:
                    lines.append("  " + str(item))
            folder_listing = "\n".join(lines)
        user_parts.append("\n【文件夹内容列表】系统已读取该目录，内容如下：\n" + (folder_listing or "(空)"))
    if context.get("code"):
        user_parts.append("\n当前代码：\n```\n" + (context["code"] or "").strip() + "\n```")
    if context.get("selected_text"):
        user_parts.append("\n选中的片段：\n```\n" + (context["selected_text"] or "").strip() + "\n```")
    if context.get("file_path"):
        user_parts.append("\n文件路径：" + str(context["file_path"]))
    if context.get("custom_command"):
        user_parts.append("\n自定义命令：" + str(context["custom_command"]))
    user_content = "\n".join(user_parts)

    # 多轮历史：context["history"] 为 [{"role":"user"|"assistant","content":"..."}, ...]，插入到 system 与当前 user 之间
    history = context.get("history") or []
    if isinstance(history, list) and history:
        history = history[-20:]  # 最多最近 20 条（约 10 轮），避免超长
        history = [
            {"role": (m.get("role") or "user").strip().lower(), "content": (m.get("content") or "").strip()}
            for m in history
            if (m.get("role") or "").strip() and (m.get("content") or "").strip()
        ]
        history = [x for x in history if x["role"] in ("user", "assistant")]
    else:
        history = []
    messages = [{"role": "system", "content": system}] + history + [{"role": "user", "content": user_content}]

    # 有文件内容或创建文件（如网站）时用接口允许的上限 8192（接口有效范围为 [1, 8192]）
    max_tokens = 8192 if (context.get("file_content") is not None or mode == "create_file") else 4096
    payload = {
        "model": _get_effective_model(model),
        "messages": messages,
        "temperature": 0.3,
        "max_tokens": max_tokens,
    }
    timeout_sec = 300
    try:
        timeout_sec = max(60, min(3600, int(os.environ.get("QWEN_REQUEST_TIMEOUT", "300"))))
    except ValueError:
        pass
    resp = _post_chat_with_retry(url, headers, payload, retries=2, timeout=timeout_sec)
    return (resp.json().get("choices") or [{}])[0].get("message", {}).get("content", "").strip()


def call_qwen_assistant_stream(
    mode: str,
    instruction: str,
    context: Optional[dict] = None,
    model: str = "qwen2.5-coder-14b",
) -> Iterator[str]:
    """
    电脑助手流式版：与 call_qwen_assistant 相同的入参和 payload，但请求 stream=True，
    逐块 yield 模型返回的 content delta。调用方拼接所有 yield 即得完整回复。
    """
    url, headers = _get_qwen_endpoint()
    context = context or {}
    file_edit_note = ""
    if context.get("file_content") is not None:
        file_edit_note = (
            " The user's local file content has been read by the system and attached below. "
            "You ONLY output the modified/polished full content—no explanation, no summary before or after. "
            "The system will reply to the user with one short sentence (e.g. 已经帮你改好了). "
            "Do NOT say you cannot access files—just output the result."
        )
    system_prompts = {
        "polish": (
            "You are a professional copy editor. The user will give you text to polish. "
            "Output ONLY the polished text in the same language, more fluent and clear. No explanations."
            + file_edit_note
        ),
        "edit_code": (
            "You are a senior programmer. The user will describe how to modify code and may provide code. "
            "Output the complete modified code in a single Markdown code block. If no code is given, output a complete script. "
            "Use the same language as the input. No extra explanation outside the block. "
            "CRITICAL: You MUST output the ENTIRE file content. Never output only a snippet or the changed part—the system will replace the whole file with your output; missing parts will be lost."
            + file_edit_note
        ),
        "complete_code": (
            "You are a code completion expert. The user will provide incomplete code. "
            "Output ONLY the completed code in a single Markdown code block. Same language and style."
        ),
        "terminal": (
            "You are a CLI assistant. The user will describe a task to do on their computer (e.g. list files, run Python, compress). "
            "Output ONLY the exact terminal command(s) to run, one per line if multiple. "
            "Use safe, common commands. Prefer Python scripts for complex tasks. No explanations."
        ),
        "folder_edit": (
            "You are a senior programmer. The user wants to modify MULTIPLE files in a folder. "
            "You are given the folder path and the content of each file. "
            "Output the modified content for each file using EXACTLY this format (no other text or explanation):\n"
            "---FILE: relative/path/filename---\n"
            "(full file content)\n"
            "---FILE: next_filename---\n"
            "(content)\n"
            "Rules: 1) Use only ---FILE: path--- blocks; 2) Path must be the same relative path as given (e.g. a.py, src/b.py); "
            "3) Output ALL files that were provided, each exactly once; 4) Apply the user's requested change to each file."
        ),
        "list_folder": (
            "You are a helpful assistant. The user asked to view the contents of a folder. "
            "You will be given the folder path and a listing of its contents (files and subfolders). "
            "Reply in Chinese with a clear, concise list or summary of what is in the folder. No need to output commands."
        ),
        "create_file": (
            "You are a senior developer. The user wants you to CREATE a new project (website, 软件/App, or Python Agent). "
            "CRITICAL: Do NOT output any design plan (设计方案), 思路, 步骤说明, or text like '让我一步步思考'. "
            "Your reply must start with ---FILE: path--- and contain ONLY ---FILE:--- blocks and optionally ---RUN:--- lines. "
            "1) Output all project files using EXACTLY this format:\n"
            "---FILE: relative/path/filename---\n"
            "(full file content)\n"
            "---FILE: next_path---\n"
            "(content)\n"
            "The system will create any folders needed (e.g. src/main.py creates src/). Paths are relative to the project folder. "
            "2) After all ---FILE:--- blocks, you MAY output one or more terminal commands for setup/build/install. Each on its own line:\n"
            "---RUN: exact_shell_command ---\n"
            "These run in order, in the created project directory. Examples: ---RUN: pip install -r requirements.txt ---, ---RUN: npm install ---, ---RUN: pio run ---, ---RUN: mkdir -p build ---. "
            "3) If the user asks for a 网站: output HTML/CSS/JS, must include index.html. "
            "4) If the user asks for 软件/App: output the project in ---FILE: path--- blocks; iOS (Swift/SwiftUI), Android (Kotlin), or cross-platform (React Native/Flutter). "
            "5) If the instruction starts with 【创造 Agent】 or is about creating an Agent: you MUST output a Python-based agent project WITH a graphical user interface (GUI). Include: main.py (or run.py) as entry point, agent logic (e.g. LLM/client calls, tool use, loop), a GUI module (use tkinter, PyQt/PySide, Gradio, Streamlit, or Flask/FastAPI with HTML/CSS/JS for web UI), requirements.txt with needed packages (e.g. openai, requests, and GUI libraries), config.py with API provider auto-detection, optional .env.example, README.md. The GUI must allow users to interact with the agent (e.g. input, output display, controls). CRITICAL: In config.py, implement automatic API provider detection based on the API key or environment variables. Support: OpenAI (base_url='https://api.openai.com/v1', models like 'gpt-3.5-turbo'), DeepSeek (base_url='https://api.deepseek.com', models like 'deepseek-chat'), DashScope/Aliyun (base_url='https://dashscope.aliyuncs.com/compatible-mode/v1', models like 'qwen-turbo'), Moonshot (base_url='https://api.moonshot.cn/v1', models like 'moonshot-v1-8k'). Detect provider from API key prefix, environment variables (DEEPSEEK_API_KEY, DASHSCOPE_API_KEY, MOONSHOT_API_KEY), or API_PROVIDER env var, then set the correct base_url and default model. Do NOT output a website or front-end only project for Agent requests. "
            "6) First line of your response must be ---FILE: or the first file path; no preamble. Do NOT show code in chat (no ```); deliver only via ---FILE:--- and ---RUN:---."
        ),
        "todo": (
            "You are a task planner. The user will describe a goal. "
            "Output a concise TODO list in this format, one per line: \"- [ ] item\" or \"- [x] item\" for done. "
            "Use Chinese. No other text."
        ),
        "plan": (
            "You are a step-by-step planner. The user will describe a task. "
            "Output a clear numbered plan (1. 2. 3. ...) in Chinese. Be concise and actionable."
        ),
        "deep_think": (
            "You are a careful reasoner and assistant. Based on the user's intent, you may:\n"
            "1) Reply with text (analysis, plan, explanation) when they ask for discussion or ideas.\n"
            "2) Create files by outputting ONLY ---FILE: relative/path--- blocks when they clearly ask you to create/make something (e.g. 做xxx网站/软件/网页/小游戏). Use format:\n"
            "---FILE: path---\n(content)\n---FILE: next---\n(content)\n"
            "The system will create any folders needed. After ---FILE:--- blocks you may add one or more ---RUN: exact_shell_command --- (one per line); they run in order in the created project directory (e.g. pip install -r requirements.txt, npm install, pio run).\n"
            "For PPT/幻灯片: use ---FILE: name.pptx --- with plain text; [IMG: url] for images. No design plan or preamble when creating—output ---FILE: blocks directly.\n"
            "3) When the user asks to run a terminal command or view system info (e.g. 通过终端查看电脑信息、执行命令), output ---RUN: exact_shell_command ---. You may output multiple ---RUN:--- lines; the system runs each and shows output. Do NOT say you cannot run the terminal.\n"
            "Use Chinese for normal replies. CRITICAL: Do NOT show code in chat (no ```). Deliver code only via ---FILE:--- and ---RUN:---; the system will save files and run commands."
        ),
        "custom_command": (
            "You are an assistant that executes user-defined commands. The user will give a command name and/or description. "
            "Interpret it and output: either a concrete action (script, command, or text result) in Chinese. "
            "If it's a script, output the full script in a Markdown code block."
        ),
    }
    system = system_prompts.get(mode) or system_prompts["deep_think"]
    user_parts = [instruction]
    if mode == "create_file":
        user_parts.append(
            "\n【必须】请直接以 ---FILE: 相对路径--- 开头输出多文件完整代码，不要输出任何设计方案、思路、步骤说明或「让我一步步思考」等文字。"
        )
        if (instruction or "").strip().startswith("【创造 Agent】") or (instruction or "").strip().startswith("【创造Agent】"):
            user_parts.append(
                "\n【创造 Agent 必须】本项目必须是基于 Python 的 Agent 项目，且必须包含图形化操作界面（GUI）。须包含：入口脚本（如 main.py 或 run.py）、Agent 核心逻辑（如调用大模型/API、工具调用、对话循环）、图形化界面模块（可使用 tkinter、PyQt/PySide、Gradio、Streamlit，或 Flask/FastAPI + HTML/CSS/JS 构建 Web 界面）、requirements.txt（需包含 GUI 相关依赖）、config.py（必须实现 API 提供商自动识别功能）、可选 .env.example、README.md。界面需支持用户与 Agent 交互（如输入框、输出显示、控制按钮等）。\n\n【API 提供商自动识别要求】在 config.py 中必须实现自动识别 API 提供商的功能：\n1. 支持识别：OpenAI（base_url='https://api.openai.com/v1'，默认模型如 'gpt-3.5-turbo'）、DeepSeek（base_url='https://api.deepseek.com'，默认模型如 'deepseek-chat'）、阿里云 DashScope（base_url='https://dashscope.aliyuncs.com/compatible-mode/v1'，默认模型如 'qwen-turbo'）、Moonshot（base_url='https://api.moonshot.cn/v1'，默认模型如 'moonshot-v1-8k'）。\n2. 识别方式：优先检查环境变量（DEEPSEEK_API_KEY、DASHSCOPE_API_KEY、MOONSHOT_API_KEY、OPENAI_API_KEY）或 API_PROVIDER 环境变量；若未设置，可根据 API key 特征或默认使用 DeepSeek。\n3. 根据识别结果自动设置正确的 base_url 和默认模型名称。\n4. 在 agent.py 中使用 OpenAI 客户端时，必须传入 base_url 参数（如 client = OpenAI(api_key=api_key, base_url=detected_base_url)）。\n不要生成纯命令行或无界面的 Agent 项目。"
            )
    if context.get("file_content") is not None:
        user_parts.append(
            "\n【重要】系统已把用户电脑上的该文件内容附在下方，你已能直接看到，请基于该内容修改并直接输出完整结果。"
            "切勿说「无法直接访问」「无法查看你的文件」等，只需输出修改后的完整文件内容。\n"
            "当前文件内容：\n```\n" + (context["file_content"] or "").strip() + "\n```"
        )
    folder_files = context.get("folder_files")
    if folder_files:
        user_parts.append("\n【说明】系统已读取文件夹内以下文件，请按用户要求修改每个文件，并用 ---FILE: 相对路径--- 格式输出全部文件。\n")
        user_parts.append("文件夹路径：" + str(context.get("folder_path", "")))
        for rel, content in folder_files:
            user_parts.append("\n--- 文件：" + rel + " ---\n" + (content or "").strip())
    folder_listing = context.get("folder_listing")
    if folder_listing is not None:
        if isinstance(folder_listing, (list, tuple)):
            lines = []
            for item in folder_listing:
                if isinstance(item, (list, tuple)):
                    name, is_file = item[0], item[1] if len(item) > 1 else True
                    lines.append("  " + name + (" (文件)" if is_file else " (文件夹)"))
                else:
                    lines.append("  " + str(item))
            folder_listing = "\n".join(lines)
        user_parts.append("\n【文件夹内容列表】系统已读取该目录，内容如下：\n" + (folder_listing or "(空)"))
    if context.get("code"):
        user_parts.append("\n当前代码：\n```\n" + (context["code"] or "").strip() + "\n```")
    if context.get("selected_text"):
        user_parts.append("\n选中的片段：\n```\n" + (context["selected_text"] or "").strip() + "\n```")
    if context.get("file_path"):
        user_parts.append("\n文件路径：" + str(context["file_path"]))
    if context.get("custom_command"):
        user_parts.append("\n自定义命令：" + str(context["custom_command"]))
    user_content = "\n".join(user_parts)
    history = context.get("history") or []
    if isinstance(history, list) and history:
        history = history[-20:]
        history = [
            {"role": (m.get("role") or "user").strip().lower(), "content": (m.get("content") or "").strip()}
            for m in history
            if (m.get("role") or "").strip() and (m.get("content") or "").strip()
        ]
        history = [x for x in history if x["role"] in ("user", "assistant")]
    else:
        history = []
    messages = [{"role": "system", "content": system}] + history + [{"role": "user", "content": user_content}]
    max_tokens = 8192 if (context.get("file_content") is not None or mode == "create_file") else 4096
    payload = {
        "model": _get_effective_model(model),
        "messages": messages,
        "temperature": 0.3,
        "max_tokens": max_tokens,
        "stream": True,
    }
    timeout_sec = 300
    try:
        timeout_sec = max(60, min(3600, int(os.environ.get("QWEN_REQUEST_TIMEOUT", "300"))))
    except ValueError:
        pass
    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=timeout_sec, stream=True)
        resp.raise_for_status()
    except Exception as e:
        # 流式请求失败时回退为非流式，一次性返回完整内容
        payload_no_stream = {k: v for k, v in payload.items() if k != "stream"}
        resp = _post_chat_with_retry(url, headers, payload_no_stream, retries=2, timeout=timeout_sec)
        full = (resp.json().get("choices") or [{}])[0].get("message", {}).get("content", "").strip()
        if full:
            yield full
        return
    for line in resp.iter_lines(decode_unicode=True):
        if not line or not line.strip().startswith("data:"):
            continue
        line = line.strip()
        if line == "data: [DONE]":
            break
        data = line[5:].strip()
        if not data:
            continue
        try:
            obj = json.loads(data)
            choices = obj.get("choices") or []
            if not choices:
                continue
            delta = choices[0].get("delta") or choices[0].get("message") or {}
            content = delta.get("content")
            if isinstance(content, str) and content:
                yield content
        except (json.JSONDecodeError, KeyError, TypeError):
            continue


def run_assistant_terminal(
    command: str,
    cwd: Optional[str] = None,
    timeout_sec: int = 60,
) -> tuple:
    """
    在指定目录下执行一条终端命令（仅允许安全命令或项目内操作）。
    返回 (ok: bool, output: str)。
    cwd 允许：项目根、桌面及其子目录（与助手可操作路径一致）。
    允许：python/pip、mkdir/cp/mv、pio、npm、ls/cat、zip、curl/wget、git/make 等，禁止 rm -rf / 等危险操作。
    """
    if not (command or "").strip():
        return False, "命令不能为空"
    root = get_project_root()
    cwd = (cwd or root).strip()
    if not cwd or not os.path.isdir(cwd):
        cwd = root
    # cwd 必须在允许的根目录之下（项目根或桌面及其子目录）
    try:
        real_cwd = os.path.realpath(cwd)
        if not _is_path_under_allowed_bases(real_cwd):
            cwd = root
    except OSError:
        cwd = root

    # 允许的命令前缀（白名单：常见开发/项目操作，禁止 rm -rf / 等）
    allowed_starts = (
        "python", "python3", "pip", "pip3",
        # 通用 / 类 Unix 命令
        "ls", "cat", "cd ", "pwd", "echo",
        "zip", "unzip", "tar ",
        "curl", "wget",
        "npx", "node ", "npm ",
        "mkdir ", "mkdir -p ", "cp ", "cp -r ", "mv ", "touch ",
        "pio ", "platformio ",
        "make ", "cmake ", "cargo ",
        "uname", "whoami", "hostname", "date", "sw_vers", "df ", "top ", "ps ",
        "lscpu", "uptime", "env", "printenv", "sysctl ", "vm_stat", "free ",
        "head ", "tail ", "wc ", "file ", "stat ", "du ", "which ", "whereis ", "find ",
        "grep ", "rg ", "diff ", "tree ",
        "ping ", "ifconfig", "ip ", "netstat ", "nslookup ", "dig ",
        "git ",
        "system_profiler", "defaults read", "launchctl list",
        # Windows 常用命令（在 Windows 上可用，在其他系统上忽略即可）
        "dir ", "type ", "cls", "chcp ",
    )
    cmd_stripped = command.strip()
    first_word = (cmd_stripped.split() or [""])[0].lower()
    if not any(cmd_stripped.lower().startswith(p) for p in allowed_starts):
        # 也允许以项目内脚本路径开头（如 python script.py）
        if first_word in ("python", "python3") and len(cmd_stripped.split()) >= 2:
            pass
        else:
            return False, f"不允许执行该命令（支持: python/pip/npx/node, ls/cat/head/tail, zip/tar, curl/wget, git/make, 系统信息 uname/df/ps/top/sw_vers 等白名单命令）。你输入的是: {cmd_stripped[:80]}"

    try:
        proc = subprocess.run(
            cmd_stripped,
            shell=True,
            cwd=cwd,
            capture_output=True,
            text=True,
            timeout=timeout_sec,
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        combined = (out + "\n" + err).strip() if err else out
        if proc.returncode != 0:
            combined = f"[exit {proc.returncode}]\n" + combined
        return proc.returncode == 0, combined or "(无输出)"
    except subprocess.TimeoutExpired:
        return False, f"命令执行超时（{timeout_sec}s）"
    except Exception as e:
        return False, str(e)


def _parse_multi_file_output(text: str) -> dict:
    """
    从模型输出中解析多文件内容。约定格式：
    ---FILE: 相对路径---
    文件内容（可多行）
    ---FILE: 下一个路径---
    ...
    返回 { "path": "content", ... }，路径统一用 / 分隔。
    若模型先输出设计方案再输出 ---FILE:---，则从第一个 ---FILE: 起截取再解析。
    """
    if not (text or "").strip():
        return {}
    lower = text.lower()
    idx = lower.find("---file:")
    if idx > 0:
        text = text[idx:]
    out: dict = {}
    pattern = re.compile(r"---FILE:\s*([^\n\-]+)---\s*\n([\s\S]*?)(?=---FILE:|---\s*FILE:|$)", re.IGNORECASE)
    for m in pattern.finditer(text):
        path = m.group(1).strip().replace("\\", "/").lstrip("/")
        content = m.group(2).rstrip()
        if path:
            out[path] = content
    if not out:
        # 兼容：整段当作 main.py
        stripped = text.strip()
        if stripped:
            out["main.py"] = stripped
    return out


def call_qwen_coder_multi_file(
    user_instruction: str, model: str = "qwen2.5-coder-14b"
) -> dict:
    """
    调用 Qwen Coder 生成多文件 MicroPython 项目。
    返回 { "main.py": "...", "lib/utils.py": "...", ... }，路径为设备上的相对路径。
    """
    url, headers = _get_qwen_endpoint()
    system_prompt = (
        "You are a senior embedded / MicroPython engineer. Target: ESP8266 (or similar) running MicroPython.\n"
        "Output a MULTI-FILE project. Use EXACTLY this format for each file (no other text):\n"
        "---FILE: relative/path/on/device.py---\n"
        "(full file content)\n"
        "---FILE: next_path.py---\n"
        "(content)\n"
        "Requirements:\n"
        "1. Must include main.py as the entry point (will be run on boot).\n"
        "2. Use only the ---FILE: path--- blocks; no Markdown fences, no explanations between files.\n"
        "3. Paths are relative to device root, e.g. main.py, lib/helper.py, config.py.\n"
        "4. For debug use print(). Keep code runnable on ESP8266 MicroPython.\n"
    )
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_instruction},
    ]
    payload = {"model": _get_effective_model(model), "messages": messages, "temperature": 0.2}
    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=240)
    content = resp.json()["choices"][0]["message"]["content"]
    return _parse_multi_file_output(content)


# =====================
# 3. 行动层：覆写 ESP8266 上的 main.py
# =====================

def write_temp_code(code: str) -> str:
    """把生成的代码写到临时文件，返回文件路径"""
    fd, path = tempfile.mkstemp(suffix=".py", prefix="esp8266_")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        f.write(code)
    return path


def build_and_upload_platformio(
    code: str,
    port: str,
    logs: List[str],
    board_id: Optional[str] = None,
    platform: Optional[str] = None,
) -> None:
    """
    使用 PlatformIO 生成临时工程、编译并（尽可能）烧录到指定串口。

    board_id / platform 可由调用方传入；未传时使用环境变量：
    - PLATFORMIO_BOARD_ID、PLATFORMIO_PLATFORM
    - PLATFORMIO: 命令名，默认 "pio"。
    """
    board_id = board_id or os.environ.get("PLATFORMIO_BOARD_ID")
    if not board_id:
        raise RuntimeError(
            "未选择开发板且未配置 PLATFORMIO_BOARD_ID。请在界面选择开发板，或在终端执行："
            " export PLATFORMIO_BOARD_ID=nodemcuv2（示例）"
        )
    platform = platform or os.environ.get("PLATFORMIO_PLATFORM", "ststm32")
    pio_cmd = os.environ.get("PLATFORMIO", "pio")

    # 统一换行，避免 \r\n 导致 replace 不匹配
    code = code.replace("\r\n", "\n").replace("\r", "\n")

    # 若前端/复用传入的是 Markdown 代码块，去掉围栏，避免把 ``` 写入 main.cpp
    code_stripped = code.strip()
    if code_stripped.startswith("```"):
        first_nl = code_stripped.find("\n")
        if first_nl != -1:
            code = code_stripped[first_nl + 1 :].rstrip()
        if code.endswith("```"):
            code = code[:-3].rstrip()
        code = code.strip()

    # ESP8266/Arduino 必须最先包含 Arduino.h；加 cstddef 可缓解 initializer_list 里 size_type 未定义
    if "Arduino.h" not in code[:800]:
        code = '#include <Arduino.h>\n#include <cstddef>\n\n' + code

    # ESP8266 的 toolchain 对部分 C++ 标准库支持有问题，去掉会触发 initializer_list 等的 include
    if platform == "espressif8266":
        for stl in ("initializer_list", "vector", "array", "list", "map", "set", "string"):
            code = re.sub(
                r'^\s*#\s*include\s*[<"]\s*' + re.escape(stl) + r'\s*[>"].*$',
                "// removed for ESP8266: #include <%s>" % stl,
                code,
                flags=re.MULTILINE | re.IGNORECASE,
            )
        # 避免隐式使用 initializer_list 的语法：将 for (x : {1,2,3}) 等写法留给模型避免，此处仅再确保无直接 include
        if "<initializer_list>" in code or "initializer_list" in code:
            code = re.sub(
                r'#\s*include\s*[<"]\s*initializer_list\s*[>"]\s*\n?',
                "// removed for ESP8266\n",
                code,
                flags=re.IGNORECASE,
            )

    # 生成代码常漏掉全局声明：只要用了 myServo 就在最前面强制加完整头部，再从原代码里删掉可能重复的声明行，确保复用/截断时也能通过编译
    if "myServo.attach" in code or "myServo.write" in code:
        required_header = (
            "#include <Arduino.h>\n"
            "#include <cstddef>\n\n"
            "#include <Servo.h>\n"
            "Servo myServo;\n"
            "const int servoPin = 5;\n\n"
        )
        body = code.lstrip()
        code = required_header + body
        # 从“原代码”部分去掉可能重复的声明行，避免重定义
        body = code[len(required_header) :]
        body = re.sub(r"^\s*Servo\s+myServo\s*;\s*\n", "", body, count=1, flags=re.MULTILINE)
        body = re.sub(r"^\s*const\s+int\s+servoPin\s*=\s*5\s*;.*\n", "", body, count=1, flags=re.MULTILINE)
        code = required_header + body

    # ESP8266/Arduino 需要 setup/loop 为 C 链接，否则链接报 undefined reference
    if 'extern "C" void setup' not in code and re.search(r'\bvoid\s+setup\s*\(\s*\)', code):
        code = re.sub(r'\bvoid\s+setup\s*\(\s*\)', 'extern "C" void setup()', code, count=1)
    if 'extern "C" void loop' not in code and re.search(r'\bvoid\s+loop\s*\(\s*\)', code):
        code = re.sub(r'\bvoid\s+loop\s*\(\s*\)', 'extern "C" void loop()', code, count=1)

    with tempfile.TemporaryDirectory(prefix="lumi_pio_") as tmpdir:
        project_dir = Path(tmpdir)
        src_dir = project_dir / "src"
        src_dir.mkdir(parents=True, exist_ok=True)
        main_cpp = src_dir / "main.cpp"
        main_cpp.write_text(code, encoding="utf-8")
        logs.append(f"[PlatformIO] 已生成 main.cpp: {main_cpp}")

        ini_path = project_dir / "platformio.ini"
        ini_lines = [
            "[env:lumi]",
            f"platform = {platform}",
            f"board = {board_id}",
            "framework = arduino",
        ]
        # 不对 ESP8266 使用 -include cstddef：该 flag 会作用到所有单元（含 .c），C 无法用 C++ 头文件，会报 cstddef: No such file or directory。main.cpp 已自带 #include <cstddef>
        ini_path.write_text("\n".join(ini_lines) + "\n", encoding="utf-8")
        logs.append(f"[PlatformIO] 开发板: {board_id}, platform: {platform}")
        logs.append(f"[PlatformIO] 已生成 platformio.ini: {ini_path}")

        # 先 clean 再编译，避免旧 .o 导致 myServo/delay 仍报未声明
        clean_cmd = [pio_cmd, "run", "-d", str(project_dir), "-t", "clean"]
        subprocess.run(clean_cmd, capture_output=True, text=True, timeout=30)
        # 编译
        compile_cmd = [pio_cmd, "run", "-d", str(project_dir)]
        logs.append(f"[PlatformIO] 编译命令: {' '.join(compile_cmd)}")
        try:
            proc = subprocess.run(
                compile_cmd,
                capture_output=True,
                text=True,
            )
        except FileNotFoundError as e:
            raise RuntimeError(
                "未找到 PlatformIO 命令（pio）。请先安装 PlatformIO，并确保 pio 在 PATH 中。"
            ) from e

        logs.append(proc.stdout)
        if proc.returncode != 0:
            logs.append(proc.stderr)
            err_tail = (proc.stderr or proc.stdout or "").strip()
            if len(err_tail) > 600:
                err_tail = "...\n" + err_tail[-600:]
            raise RuntimeError(
                "PlatformIO 编译失败。错误摘要：\n" + (err_tail or "无输出")
            )

        # 烧录
        upload_cmd = [pio_cmd, "run", "-d", str(project_dir), "-t", "upload"]
        if port:
            upload_cmd += ["--upload-port", port]
        logs.append(f"[PlatformIO] 烧录命令: {' '.join(upload_cmd)}")

        proc2 = subprocess.run(
            upload_cmd,
            capture_output=True,
            text=True,
        )
        logs.append(proc2.stdout)
        if proc2.returncode != 0:
            logs.append(proc2.stderr)
            raise RuntimeError("PlatformIO 烧录失败，请查看上方烧录日志。")

        logs.append("[PlatformIO] 编译并烧录完成。")


def ping_qwen_model(max_timeout: int = 10) -> dict:
    """
    轻量级探测当前 Qwen / 本地模型接口是否可用，用于在 UI 中展示连通性和简单延迟。
    不依赖业务 prompt，只发送一个极短的请求并读取 1 个 token。
    """
    t0 = time.monotonic()
    try:
        url, headers = _get_qwen_endpoint()
    except Exception as e:
        dt = (time.monotonic() - t0) * 1000
        return {
            "ok": False,
            "stage": "config",
            "latency_ms": int(dt),
            "error": str(e),
        }

    payload = {
        "model": _get_effective_model("qwen2.5-coder-14b"),
        "messages": [
            {"role": "system", "content": "You are a lightweight health check."},
            {"role": "user", "content": "ping"},
        ],
        "temperature": 0.0,
        "max_tokens": 1,
    }

    try:
        resp = _post_chat_with_retry(url, headers, payload, retries=0, timeout=max_timeout)
        _ = resp.json()
        dt = (time.monotonic() - t0) * 1000
        return {
            "ok": True,
            "stage": "request",
            "latency_ms": int(dt),
        }
    except Exception as e:
        dt = (time.monotonic() - t0) * 1000
        return {
            "ok": False,
            "stage": "request",
            "latency_ms": int(dt),
            "error": str(e),
        }


def clear_lumi_cache(logs: Optional[List[str]] = None) -> None:
    """
    清理 Lumi 生成的临时 / 缓存文件。

    当前策略：
    - 扫描系统临时目录（tempfile.gettempdir()）下名称以以下前缀开头的条目：
      - esp8266_
      - lumi_pio_
      - lumi_arduino_
    - 如为目录则递归删除，如为文件则直接删除。
    """
    if logs is None:
        logs = []

    tmp_root = Path(tempfile.gettempdir())
    prefixes = ("esp8266_", "lumi_pio_", "lumi_arduino_")
    removed = 0

    for child in tmp_root.iterdir():
        if not any(child.name.startswith(p) for p in prefixes):
            continue
        try:
            if child.is_dir():
                shutil.rmtree(child, ignore_errors=True)
            else:
                try:
                    child.unlink()
                except FileNotFoundError:
                    pass
            removed += 1
        except Exception as e:
            logs.append(f"删除缓存失败: {child} -> {e}")

    logs.append(f"已清理 {removed} 个 Lumi 缓存条目。")


def flash_micropython_main(port: str, src_path: str):
    """
    使用 mpremote 把 src_path 上传为板子上的 main.py。
    在实际拷贝前会先简单检查一次指定串口上是否运行着 MicroPython，
    如果检测失败，会抛出更友好的错误而不是直接返回晦涩的系统异常。
    """
    # 先做一次简单探测，避免在非 MicroPython 设备上直接尝试 cp
    check_cmd = [
        "mpremote",
        "connect",
        port,
        "exec",
        "import sys; print('MPY')",
    ]
    print("检查 MicroPython 设备:", " ".join(check_cmd))
    try:
        probe = subprocess.run(
            check_cmd,
            check=True,
            capture_output=True,
            text=True,
            timeout=10,
        )
    except FileNotFoundError as e:
        raise RuntimeError("未找到 mpremote 工具，请先通过 pip 安装 mpremote。") from e
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"无法在端口 {port} 上检测到 MicroPython 设备，请确认已正确刷入 MicroPython 固件。"
        ) from e
    except subprocess.TimeoutExpired as e:
        raise RuntimeError(
            f"检查端口 {port} 上的 MicroPython 状态超时，请确认设备已连接并重试。"
        ) from e

    if "MPY" not in (probe.stdout or ""):
        raise RuntimeError(
            f"端口 {port} 未返回预期的 MicroPython 响应，请确认该设备已刷入 MicroPython 固件。"
        )

    # 真正执行拷贝
    cmd = [
        "mpremote",
        "connect",
        port,
        "cp",
        src_path,
        ":main.py",
    ]
    print("执行命令:", " ".join(cmd))
    try:
        subprocess.run(cmd, check=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"通过 mpremote 向端口 {port} 烧录 main.py 失败，"
            f"请检查连接是否稳定或设备是否正在被其他程序占用。（底层错误: {e}）"
        ) from e

    print("上传完成，重启板子后会自动运行 main.py")


def flash_micropython_files(
    port: str, files_dict: dict, logs: Optional[List[str]] = None
) -> None:
    """
    将多文件项目上传到 MicroPython 设备。
    files_dict: { "main.py": "content", "lib/foo.py": "content", ... }
    会先检测设备是否为 MicroPython，再按路径创建目录并拷贝每个文件。
    """
    if logs is None:
        logs = []

    def log(msg: str):
        logs.append(msg)
        print(msg)

    # 设备检测（与 flash_micropython_main 一致）
    check_cmd = ["mpremote", "connect", port, "exec", "import sys; print('MPY')"]
    try:
        probe = subprocess.run(
            check_cmd, check=True, capture_output=True, text=True, timeout=10
        )
    except FileNotFoundError as e:
        raise RuntimeError("未找到 mpremote，请先通过 pip 安装 mpremote。") from e
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"无法在端口 {port} 上检测到 MicroPython 设备，请确认已正确刷入 MicroPython 固件。"
        ) from e
    if "MPY" not in (probe.stdout or ""):
        raise RuntimeError(
            f"端口 {port} 未返回预期的 MicroPython 响应，请确认该设备已刷入 MicroPython 固件。"
        )

    # 收集需要创建的远程目录（路径用 /）
    remote_dirs: set = set()
    for rel_path in files_dict:
        parts = rel_path.replace("\\", "/").split("/")
        if len(parts) > 1:
            for i in range(1, len(parts)):
                remote_dirs.add("/".join(parts[:i]))

    with tempfile.TemporaryDirectory(prefix="lumi_mpy_") as tmpdir:
        root = Path(tmpdir)
        for rel_path, content in files_dict.items():
            rel_path = rel_path.replace("\\", "/").lstrip("/")
            if not rel_path:
                continue
            local_path = root / rel_path
            local_path.parent.mkdir(parents=True, exist_ok=True)
            local_path.write_text(content, encoding="utf-8")

        # 在设备上创建目录（mpremote fs mkdir）
        for d in sorted(remote_dirs):
            mkdir_cmd = ["mpremote", "connect", port, "fs", "mkdir", ":" + d]
            try:
                subprocess.run(mkdir_cmd, capture_output=True, text=True, timeout=10)
            except Exception:
                pass  # 已存在会报错，忽略

        for rel_path in files_dict:
            rel_path = rel_path.replace("\\", "/").lstrip("/")
            if not rel_path:
                continue
            local_path = root / rel_path
            if not local_path.is_file():
                continue
            remote = ":" + rel_path
            cmd = ["mpremote", "connect", port, "cp", str(local_path), remote]
            log(f"上传 {rel_path} -> 设备 {remote}")
            try:
                subprocess.run(cmd, check=True, timeout=30)
            except subprocess.CalledProcessError as e:
                raise RuntimeError(
                    f"上传 {rel_path} 到设备失败: {e}"
                ) from e

    log("多文件上传完成，请重启或复位设备。")


def get_project_root() -> str:
    """返回项目根目录，用于解析相对路径。默认 ~/Desktop，可通过 LUMI_PROJECT_ROOT 覆盖。"""
    root = os.environ.get("LUMI_PROJECT_ROOT", "").strip()
    if not root:
        root = os.path.expanduser("~/Desktop")
    return os.path.abspath(root)


def _resolve_editable_path(relative_path: str) -> str:
    """将相对路径解析为绝对路径，并校验不跳出项目根。"""
    root = get_project_root()
    abs_path = os.path.normpath(os.path.join(root, relative_path))
    try:
        real_root = os.path.realpath(root)
        real_path = os.path.realpath(abs_path)
    except OSError:
        real_root, real_path = root, abs_path
    if real_path != real_root and not real_path.startswith(real_root + os.sep):
        raise ValueError("路径不能超出项目根目录")
    return abs_path


def edit_file_preview(
    relative_path: str,
    user_instruction: str,
    selected_text: Optional[str] = None,
    context_files: Optional[List[dict]] = None,
) -> str:
    """
    AI 预览编辑结果：根据指令与上下文生成新内容，不写回磁盘。
    relative_path: 相对于项目根的路径。
    context_files: [{"path": "rel/path", "content": "..."}]，若只有 path 则从磁盘读取（在项目根下）。
    返回：生成的新文件内容。
    """
    abs_path = _resolve_editable_path(relative_path)
    if not os.path.exists(abs_path):
        raise FileNotFoundError(f"文件不存在: {abs_path}")

    with open(abs_path, "r", encoding="utf-8") as f:
        old_content = f.read()

    ctx_with_content: List[dict] = []
    if context_files:
        root = get_project_root()
        for ctx in context_files:
            path = (ctx.get("path") or "").strip()
            content = ctx.get("content")
            if content is not None:
                ctx_with_content.append({"path": path, "content": content})
            elif path:
                try:
                    p_abs = os.path.normpath(os.path.join(root, path))
                    if p_abs.startswith(root + os.sep) and os.path.isfile(p_abs):
                        with open(p_abs, "r", encoding="utf-8") as f2:
                            ctx_with_content.append({"path": path, "content": f2.read()})
                except Exception:
                    pass

    return call_qwen_file_editor(
        abs_path,
        old_content,
        user_instruction,
        selected_text=selected_text,
        context_files=ctx_with_content or None,
    )


def edit_file_apply(relative_path: str, new_content: str) -> None:
    """将 AI 生成的内容写回指定文件（路径必须在项目根下）。"""
    abs_path = _resolve_editable_path(relative_path)
    Path(abs_path).write_text(new_content, encoding="utf-8")


def edit_desktop_file(relative_path: str, user_instruction: str) -> str:
    """
    使用 Qwen Coder 2.5 修改项目根目录下的一个文件（默认桌面）。
    relative_path: 相对于项目根的路径，例如 'test.py' 或 'project/notes.md'
    返回：修改后的完整文件内容（已写回磁盘）。
    """
    new_content = edit_file_preview(relative_path, user_instruction)
    edit_file_apply(relative_path, new_content)
    return new_content


def _to_github_raw_url(url: str) -> str:
    """
    将常见的 GitHub 仓库文件地址转换为 raw 地址。
    支持两种形式：
    - https://github.com/user/repo/blob/branch/path/to/file.py
    - https://raw.githubusercontent.com/user/repo/branch/path/to/file.py （原样返回）
    其他 URL 则原样返回。
    """
    if "raw.githubusercontent.com" in url:
        return url
    if "github.com" in url and "/blob/" in url:
        return url.replace("github.com", "raw.githubusercontent.com").replace("/blob/", "/")
    return url


def download_github_file(url: str) -> str:
    """
    从 GitHub（或任意可直接访问的原始文本 URL）下载文件内容，返回字符串。
    主要用于把开源示例代码一键烧录到 ESP8266。
    """
    if not url:
        raise ValueError("URL 不能为空")

    raw_url = _to_github_raw_url(url.strip())
    resp = requests.get(raw_url, timeout=60)
    resp.raise_for_status()
    # 假定是文本文件（MicroPython / Python / 配置等）
    return resp.text


# =====================
# GitHub 搜索并下载（根据用户自然语言指令）
# =====================

GITHUB_API = "https://api.github.com"


def _github_request(path: str, params: Optional[dict] = None) -> dict:
    """发起 GitHub API 请求（未鉴权 60 次/小时）。"""
    url = GITHUB_API + path if path.startswith("/") else GITHUB_API + "/" + path
    headers = {"Accept": "application/vnd.github.v3+json"}
    token = os.environ.get("GITHUB_TOKEN")
    if token:
        headers["Authorization"] = f"Bearer {token}"
    resp = requests.get(url, params=params, headers=headers, timeout=15)
    resp.raise_for_status()
    return resp.json()


def search_github_repositories(q: str, per_page: int = 5) -> List[dict]:
    """
    按关键词搜索 GitHub 仓库。返回列表，每项含 full_name, default_branch, html_url。
    """
    data = _github_request(
        "/search/repositories",
        params={"q": q, "per_page": per_page, "sort": "stars"},
    )
    items = data.get("items") or []
    return [
        {
            "full_name": r.get("full_name"),
            "default_branch": r.get("default_branch", "main"),
            "html_url": r.get("html_url"),
        }
        for r in items
        if r.get("full_name")
    ]


def get_repo_main_file_url(repo_full_name: str, branch: Optional[str] = None) -> Optional[str]:
    """
    获取仓库根目录下的主入口文件（main.py 或第一个 .py 文件）的 raw 下载地址。
    若未指定 branch 则先请求仓库信息取 default_branch。
    """
    parts = repo_full_name.split("/", 1)
    if len(parts) != 2:
        return None
    owner, repo = parts
    if not branch:
        try:
            repo_info = _github_request(f"/repos/{owner}/{repo}")
            branch = repo_info.get("default_branch") or "main"
        except Exception:
            branch = "main"
    try:
        contents = _github_request(f"/repos/{owner}/{repo}/contents", params={"ref": branch})
    except Exception:
        return None
    if not isinstance(contents, list):
        return None
    # 优先 main.py，否则第一个 .py 文件
    candidates = [c for c in contents if c.get("type") == "file" and (c.get("name") or "").endswith(".py")]
    main_py = next((c for c in candidates if (c.get("name") or "").lower() == "main.py"), None)
    entry = main_py or (candidates[0] if candidates else None)
    if not entry:
        return None
    return entry.get("download_url")


def _call_qwen_github_resolve(instruction: str) -> str:
    """
    用模型解析用户指令，返回一行：要么是 GitHub raw 文件 URL，要么是 "SEARCH 英文搜索词"。
    """
    url, headers = _get_qwen_endpoint()
    system = (
        "你只回复一行文字。用户希望从 GitHub 下载代码并烧录到设备。"
        "若你知道一个合适的、可直接下载的 GitHub 上的示例/飞控 main.py 的 raw 地址（必须是 https://raw.githubusercontent.com/... 形式），直接输出该 URL。"
        "否则输出：SEARCH 加一个简短的英文搜索词（用于 GitHub 仓库搜索），例如 SEARCH esp8266 micropython flight controller。不要解释、不要换行。"
    )
    payload = {
        "model": _get_effective_model("qwen2.5-coder-14b"),
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": instruction},
        ],
        "temperature": 0.1,
        "max_tokens": 256,
    }
    resp = _post_chat_with_retry(url, headers, payload, retries=1, timeout=60)
    data = resp.json()
    content = (data.get("choices") or [{}])[0].get("message") or {}
    line = (content.get("content") or "").strip().split("\n")[0].strip()
    return line


def github_search_and_download(instruction: str, logs: Optional[List[str]] = None) -> Tuple[str, str]:
    """
    根据用户自然语言指令（如「帮我从 GitHub 下载完整飞控代码并烧录到设备」），
    自动解析为 GitHub URL 或搜索词，搜索/下载源文件内容。返回 (code_content, source_info)。
    source_info 用于日志展示，如仓库名或 URL。
    """
    if logs is None:
        logs = []

    def log(msg: str):
        logs.append(msg)
        print(msg)

    log("正在解析指令并确定 GitHub 来源…")
    try:
        line = _call_qwen_github_resolve(instruction)
    except Exception as e:
        log(f"模型解析失败: {e}，将使用默认搜索词。")
        line = "SEARCH esp8266 micropython flight controller"

    line = (line or "").strip()
    if not line:
        line = "SEARCH esp8266 micropython flight controller"

    if line.upper().startswith("SEARCH "):
        query = line[7:].strip() or "esp8266 micropython flight controller"
        log(f"使用 GitHub 搜索: {query}")
        repos = search_github_repositories(query, per_page=5)
        if not repos:
            raise RuntimeError("GitHub 未找到匹配的仓库，请换一个描述或稍后重试。")
        for r in repos:
            full_name = r.get("full_name")
            branch = r.get("default_branch") or "main"
            if not full_name:
                continue
            log(f"尝试仓库: {full_name}")
            download_url = get_repo_main_file_url(full_name, branch)
            if not download_url:
                continue
            try:
                code = download_github_file(download_url)
                if code and len(code.strip()) > 50:
                    log(f"已从 {full_name} 下载 main.py")
                    return code, full_name
            except Exception as e:
                log(f"下载 {download_url} 失败: {e}")
                continue
        raise RuntimeError("在已搜索到的仓库中未找到可用的 main.py，请尝试提供具体 GitHub 链接。")
    else:
        url = line
        if "raw.githubusercontent.com" not in url and "github.com" in url and "/blob/" in url:
            url = _to_github_raw_url(url)
        elif "raw.githubusercontent.com" not in url:
            url = url.strip()
        log(f"正在下载: {url}")
        code = download_github_file(url)
        return code, url


# =====================
# 自动化脚本工具箱
# =====================

def probe_micropython(port: str) -> tuple:
    """检测指定串口是否为 MicroPython 设备。返回 (ok: bool, message: str)"""
    if not port:
        return False, "未指定串口"
    try:
        check_cmd = [
            "mpremote",
            "connect",
            port,
            "exec",
            "import sys; print('MPY')",
        ]
        proc = subprocess.run(
            check_cmd,
            capture_output=True,
            text=True,
            timeout=10,
        )
        if proc.returncode == 0 and "MPY" in (proc.stdout or ""):
            return True, f"端口 {port} 已检测到 MicroPython。"
        return False, f"端口 {port} 未返回 MicroPython 响应，请确认已刷入 MicroPython 固件。"
    except FileNotFoundError:
        return False, "未找到 mpremote，请先通过 pip 安装 mpremote。"
    except subprocess.TimeoutExpired:
        return False, f"检测 {port} 超时，请检查连接。"
    except Exception as e:
        return False, str(e)


def check_platformio_env() -> tuple:
    """检查 PlatformIO 是否可用。返回 (ok: bool, message: str, version: Optional[str])"""
    pio_cmd = os.environ.get("PLATFORMIO", "pio")
    path = shutil.which(pio_cmd)
    if not path:
        return False, "未找到 PlatformIO（pio），请先安装并加入 PATH。", None
    try:
        proc = subprocess.run(
            [pio_cmd, "--version"],
            capture_output=True,
            text=True,
            timeout=10,
        )
        out = (proc.stdout or proc.stderr or "").strip()
        if proc.returncode == 0:
            return True, "PlatformIO 环境正常。", out or None
        return False, out or "执行 pio --version 失败。", None
    except Exception as e:
        return False, str(e), None


def install_missing_dependencies(logs: Optional[List[str]] = None) -> Tuple[bool, str]:
    """
    检测并自动安装项目所需依赖：先执行 pip install -r requirements.txt，
    再安装可选包 openai、ruff、platformio。需网络权限。
    requirements.txt 取自本模块所在目录（即项目/代码根目录），与 LUMI_PROJECT_ROOT 无关。
    返回 (成功与否, 摘要信息)。
    """
    def log(msg: str):
        if logs is not None:
            logs.append(msg)
        print(msg)

    # 项目/代码根目录：usb_iot_agent.py 所在目录（此处有 requirements.txt）
    root = os.path.dirname(os.path.abspath(__file__))
    req_path = os.path.join(root, "requirements.txt")
    if not os.path.isfile(req_path):
        log("未找到 requirements.txt，跳过依赖安装。")
        return False, "项目根目录下无 requirements.txt"

    log("正在检测并安装依赖（需联网）…")
    ok_all = True
    summary = []

    # 1. 安装 requirements.txt 中的全部包
    try:
        log("[1/2] 执行: pip install -r requirements.txt")
        proc = subprocess.run(
            [sys.executable, "-m", "pip", "install", "-r", req_path],
            capture_output=True,
            text=True,
            timeout=300,
            cwd=root,
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        if out:
            log(out)
        if err:
            log(err)
        if proc.returncode != 0:
            ok_all = False
            summary.append("requirements.txt 安装失败")
        else:
            summary.append("requirements.txt 已安装/更新")
    except subprocess.TimeoutExpired:
        log("pip install -r requirements.txt 超时（300s）")
        ok_all = False
        summary.append("安装超时")
    except Exception as e:
        log("执行 pip 失败: %s" % e)
        ok_all = False
        summary.append(str(e))

    # 2. 可选包：openai, ruff, platformio
    optional = ["openai", "ruff", "platformio"]
    try:
        log("[2/2] 可选包: %s" % ", ".join(optional))
        proc = subprocess.run(
            [sys.executable, "-m", "pip", "install", "-q"] + optional,
            capture_output=True,
            text=True,
            timeout=120,
            cwd=root,
        )
        if proc.returncode != 0 and proc.stderr:
            log(proc.stderr.strip())
        if proc.returncode == 0:
            summary.append("可选包已安装/跳过（已存在）")
        else:
            summary.append("部分可选包安装失败（可忽略）")
    except Exception as e:
        log("可选包安装异常: %s" % e)
        summary.append("可选包: %s" % e)

    msg = "; ".join(summary)
    log("完成: %s" % msg)
    return ok_all, msg


def export_code_to_desktop(code: str, suffix: str = ".py") -> str:
    """将代码导出到桌面，返回写入的绝对路径。"""
    desktop = Path(os.path.expanduser("~/Desktop"))
    desktop.mkdir(parents=True, exist_ok=True)
    stamp = time.strftime("%Y%m%d_%H%M%S")
    name = f"lumi_export_{stamp}{suffix}"
    path = desktop / name
    path.write_text(code, encoding="utf-8")
    return str(path)


def crawl_and_download(
    url: str,
    download_path: str,
    logs: Optional[List[str]] = None,
    timeout_sec: int = 120,
) -> tuple:
    """
    使用 requests 爬取指定 URL 并保存到本地。
    download_path: 可为目录（则根据 URL 或 Content-Disposition 生成文件名）或完整文件路径。
    若为空则默认保存到桌面。
    返回 (ok, saved_path 或 error_message)。
    """
    if logs is None:
        logs = []

    def log(msg: str):
        logs.append(msg)
        print(msg)

    url = (url or "").strip()
    if not url:
        return False, "请输入要爬取的网址"
    if not url.startswith(("http://", "https://")):
        return False, "网址需以 http:// 或 https:// 开头"

    base_dir = os.path.expanduser("~/Desktop")
    if (download_path or "").strip():
        raw = os.path.expanduser(download_path.strip())
        if os.path.isdir(raw):
            base_dir = raw
            save_path = None
        else:
            parent = os.path.dirname(raw)
            if parent:
                Path(parent).mkdir(parents=True, exist_ok=True)
            save_path = os.path.abspath(raw)
    else:
        save_path = None

    try:
        log(f"正在请求: {url}")
        resp = requests.get(
            url,
            timeout=timeout_sec,
            stream=True,
            headers={"User-Agent": "Mozilla/5.0 (compatible; LumiCrawler/1.0)"},
        )
        resp.raise_for_status()
    except requests.exceptions.Timeout:
        return False, f"请求超时（{timeout_sec}s）"
    except requests.exceptions.RequestException as e:
        return False, f"请求失败: {e}"
    except Exception as e:
        return False, str(e)

    if save_path is None:
        filename = None
        cd = resp.headers.get("Content-Disposition")
        if cd and "filename=" in cd:
            m = re.search(r'filename[*]?=(?:UTF-8\'\')?["\']?([^"\';]+)', cd, re.I)
            if m:
                filename = unquote(m.group(1).strip().strip('"\''))
        if not filename:
            parsed = urlparse(url)
            path_part = unquote(parsed.path or "")
            filename = os.path.basename(path_part) or "index.html"
        filename = re.sub(r'[<>:"/\\|?*]', "_", filename)
        Path(base_dir).mkdir(parents=True, exist_ok=True)
        save_path = os.path.join(base_dir, filename)

    content_type = (resp.headers.get("Content-Type") or "").lower()
    is_text = "text/" in content_type or "json" in content_type or "javascript" in content_type
    log(f"保存到: {save_path} ({'文本' if is_text else '二进制'})")

    try:
        if is_text:
            text = resp.content.decode(resp.encoding or "utf-8", errors="replace")
            Path(save_path).write_text(text, encoding="utf-8")
        else:
            with open(save_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=65536):
                    if chunk:
                        f.write(chunk)
        size = os.path.getsize(save_path)
        log(f"已下载完成，大小: {size} 字节")
        return True, save_path
    except OSError as e:
        return False, f"写入文件失败: {e}"


def open_project_root_in_explorer() -> str:
    """在系统文件管理器中打开项目根目录。返回项目根路径。"""
    root = get_project_root()
    if sys.platform == "darwin":
        subprocess.run(["open", root], check=True, timeout=5)
    elif sys.platform == "win32":
        subprocess.run(["explorer", os.path.normpath(root)], check=True, timeout=5)
    else:
        subprocess.run(["xdg-open", root], check=True, timeout=5)
    return root


def read_file_for_preview(path: str) -> Tuple[bool, str, str]:
    """读取文件内容供前端预览。仅允许桌面或项目根下的文件。若 path 是目录则尝试读取其下 index.html。返回 (ok, content, error_message)。"""
    path = os.path.normpath(os.path.expanduser(path.strip()))
    if not _is_path_under_allowed_bases(path):
        return False, "", "仅允许预览桌面或项目根下的文件"
    if os.path.isdir(path):
        index_path = os.path.join(path, "index.html")
        if os.path.isfile(index_path):
            path = index_path
        else:
            return False, "", "该路径是文件夹且其中没有 index.html，请使用「网页预览」在页面内运行"
    if not os.path.isfile(path):
        return False, "", "文件不存在"
    ok, content, err = read_file_content_for_assistant(path)
    return ok, content or "", err or ""


def open_file_in_system(path: str) -> Tuple[bool, str]:
    """用系统默认应用打开文件。仅允许桌面或项目根下的路径。返回 (ok, error_message)。"""
    path = os.path.normpath(os.path.expanduser(path.strip()))
    if not _is_path_under_allowed_bases(path):
        return False, "仅允许打开桌面或项目根下的文件"
    if not os.path.isfile(path):
        return False, "文件不存在"
    try:
        if sys.platform == "darwin":
            subprocess.run(["open", path], check=True, timeout=5)
        elif sys.platform == "win32":
            os.startfile(path)
        else:
            subprocess.run(["xdg-open", path], check=True, timeout=5)
        return True, ""
    except Exception as e:
        return False, str(e)


def open_folder_in_system(path: str) -> Tuple[bool, str]:
    """在系统文件管理器中打开目录；若 path 为文件则打开其所在目录。仅允许桌面或项目根下。返回 (ok, error_message)。"""
    path = os.path.normpath(os.path.expanduser(path.strip()))
    if not _is_path_under_allowed_bases(path):
        return False, "仅允许打开桌面或项目根下的路径"
    dir_path = path if os.path.isdir(path) else os.path.dirname(path)
    if not os.path.isdir(dir_path):
        return False, "目录不存在"
    try:
        if sys.platform == "darwin":
            subprocess.run(["open", dir_path], check=True, timeout=5)
        elif sys.platform == "win32":
            subprocess.run(["explorer", os.path.normpath(dir_path)], check=True, timeout=5)
        else:
            subprocess.run(["xdg-open", dir_path], check=True, timeout=5)
        return True, ""
    except Exception as e:
        return False, str(e)


def open_xcode_project(path: str) -> Tuple[bool, str]:
    """用 Xcode 打开工程。path 为项目目录或 .xcodeproj 路径；仅允许桌面或项目根下。返回 (ok, error_message)。"""
    path = os.path.normpath(os.path.expanduser(path.strip()))
    if not _is_path_under_allowed_bases(path):
        return False, "仅允许打开桌面或项目根下的路径"
    xcodeproj_path = None
    if path.endswith(".xcodeproj") and os.path.isdir(path):
        xcodeproj_path = path
    else:
        dir_path = path if os.path.isdir(path) else os.path.dirname(path)
        if not os.path.isdir(dir_path):
            return False, "目录不存在"
        try:
            for name in os.listdir(dir_path):
                if name.endswith(".xcodeproj") and os.path.isdir(os.path.join(dir_path, name)):
                    xcodeproj_path = os.path.join(dir_path, name)
                    break
        except OSError:
            pass
    if not xcodeproj_path:
        return False, "该路径下未找到 .xcodeproj，不是 Xcode 工程"
    try:
        if sys.platform == "darwin":
            subprocess.run(["open", xcodeproj_path], check=True, timeout=5)
            return True, ""
        return False, "仅支持在 macOS 上用 Xcode 打开"
    except Exception as e:
        return False, str(e)


def list_device_files(port: str) -> tuple:
    """列出 MicroPython 设备根目录下的文件。返回 (ok, message 或文件列表文本)。"""
    if not port or not port.strip():
        return False, "未指定串口"
    port = port.strip()
    try:
        cmd = ["mpremote", "connect", port, "fs", "ls", "/"]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=15)
        if proc.returncode != 0:
            return False, (proc.stderr or proc.stdout or "列出文件失败").strip()
        return True, (proc.stdout or "").strip() or "(根目录为空)"
    except FileNotFoundError:
        return False, "未找到 mpremote，请先通过 pip 安装 mpremote。"
    except subprocess.TimeoutExpired:
        return False, f"连接 {port} 超时，请检查设备。"
    except Exception as e:
        return False, str(e)


def soft_reset_device(port: str) -> tuple:
    """软复位 MicroPython 设备。返回 (ok, message)。"""
    if not port or not port.strip():
        return False, "未指定串口"
    port = port.strip()
    try:
        cmd = [
            "mpremote",
            "connect",
            port,
            "exec",
            "import machine; machine.soft_reset()",
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
        if proc.returncode != 0:
            return False, (proc.stderr or proc.stdout or "软复位失败").strip()
        return True, f"已向 {port} 发送软复位。"
    except FileNotFoundError:
        return False, "未找到 mpremote，请先通过 pip 安装 mpremote。"
    except subprocess.TimeoutExpired:
        return True, f"已发送软复位（{port} 可能已重启）。"
    except Exception as e:
        return False, str(e)


def mip_install_on_device(port: str, package: str, logs: Optional[List[str]] = None) -> tuple:
    """在 MicroPython 设备上通过 mip 安装包。返回 (ok, message)。"""
    if logs is None:
        logs = []
    port = (port or "").strip()
    package = (package or "").strip()
    if not port:
        return False, "未指定串口"
    if not package:
        return False, "未指定要安装的包名（如 umqtt.simple、aioble）"
    try:
        cmd = ["mpremote", "connect", port, "mip", "install", package]
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        for line in (out + "\n" + err).strip().splitlines():
            if line.strip():
                logs.append(line)
        if proc.returncode != 0:
            return False, err or out or "mip 安装失败"
        return True, f"已在设备 {port} 上安装 {package}。"
    except FileNotFoundError:
        return False, "未找到 mpremote，请先通过 pip 安装 mpremote。"
    except subprocess.TimeoutExpired:
        return False, "安装超时，请检查网络或稍后重试。"
    except Exception as e:
        return False, str(e)


def read_device_repl(
    port: str, duration_sec: int = 8, logs: Optional[List[str]] = None
) -> tuple:
    """
    连接设备并捕获一段时间内的 REPL 输出（设备上正在运行的程序或 REPL 的打印）。
    通过 exec 执行一段会持续打印的代码并收集输出。duration_sec 建议 5–30。
    返回 (ok, output_text)。
    """
    if logs is None:
        logs = []
    port = (port or "").strip()
    if not port:
        return False, "未指定串口"
    duration_sec = max(2, min(60, int(duration_sec)))
    # 在设备上执行一段会周期性打印的代码，便于用户看到设备状态
    script = (
        "import time; start=time.time(); "
        "print('--- 设备 REPL 输出 (每 1s 采样) ---'); "
        "import gc; "
        "while time.time()-start < %d: print('mem_free:', gc.mem_free()); time.sleep(1)"
        % duration_sec
    )
    try:
        cmd = ["mpremote", "connect", port, "exec", script]
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=duration_sec + 15
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        combined = (out + "\n" + err).strip()
        for line in combined.splitlines():
            if line.strip():
                logs.append(line)
        return True, combined or "(无输出)"
    except FileNotFoundError:
        return False, "未找到 mpremote，请先通过 pip 安装 mpremote。"
    except subprocess.TimeoutExpired:
        return False, "连接或读取设备超时。"
    except Exception as e:
        return False, str(e)


def run_project_script(
    relative_path: str, logs: Optional[List[str]] = None
) -> tuple:
    """在项目根目录下执行指定 Python 脚本。返回 (ok, output)。"""
    if logs is None:
        logs = []
    root = get_project_root()
    rel = (relative_path or "").strip()
    path = os.path.normpath(os.path.join(root, rel))
    real_path = os.path.abspath(path)
    real_root = os.path.abspath(root)
    if real_path != real_root and not real_path.startswith(real_root + os.sep):
        return False, "路径必须在项目根目录下。"
    if not os.path.isfile(path):
        return False, f"文件不存在: {path}"
    exe = "python3" if shutil.which("python3") else "python"
    try:
        proc = subprocess.run(
            [exe, path],
            cwd=root,
            capture_output=True,
            text=True,
            timeout=60,
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        for line in (out + "\n" + err).splitlines():
            logs.append(line)
        if proc.returncode != 0:
            return False, err or out or "脚本执行失败"
        return True, out or "(无标准输出)"
    except FileNotFoundError:
        return False, f"未找到 {exe}，请确保已安装 Python。"
    except subprocess.TimeoutExpired:
        return False, "脚本执行超时（60s）。"
    except Exception as e:
        return False, str(e)


def ruff_check_project(
    glob_pattern: Optional[str] = None, logs: Optional[List[str]] = None
) -> tuple:
    """在项目根目录下运行 ruff check。返回 (ok, output)。"""
    if logs is None:
        logs = []
    root = get_project_root()
    cmd = [shutil.which("ruff") or "ruff", "check", root]
    if glob_pattern and glob_pattern.strip():
        cmd.extend(["--glob", glob_pattern.strip()])
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=30
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        for line in (out + "\n" + err).splitlines():
            logs.append(line)
        if proc.returncode != 0:
            return True, out or err or "ruff 发现部分问题（见上方）。"
        return True, out or "未发现问题。"
    except FileNotFoundError:
        return False, "未找到 ruff，请先安装: pip install ruff"
    except Exception as e:
        return False, str(e)


def check_python_env() -> tuple:
    """检查本机 Python 版本。返回 (ok, message, version_str)。"""
    exe = "python3" if shutil.which("python3") else "python"
    try:
        proc = subprocess.run(
            [exe, "--version"],
            capture_output=True,
            text=True,
            timeout=5,
        )
        out = (proc.stdout or proc.stderr or "").strip()
        if proc.returncode == 0 and out:
            return True, "本机 Python 环境正常。", out
        return False, out or "无法获取版本。", None
    except FileNotFoundError:
        return False, f"未找到 {exe}，请确保已安装 Python 并加入 PATH。", None
    except Exception as e:
        return False, str(e), None


def search_in_project(
    query: str,
    glob_pattern: Optional[str],
    logs: List[str],
    max_matches: int = 200,
) -> tuple:
    """
    在项目根目录下按关键字搜索文件内容的简单工具。
    glob_pattern: 可选，逗号分隔的通配符模式，如 \"*.py,*.js\"。
    """
    root = Path(get_project_root())
    query = (query or "").strip()
    if not query:
        return False, {"error": "搜索关键字不能为空"}

    ignore_dirs = {".git", ".venv", "__pycache__", "node_modules", ".cursor"}
    patterns: Optional[List[str]] = None
    if glob_pattern:
        patterns = [p.strip() for p in glob_pattern.split(",") if p.strip()]

    matches: List[str] = []

    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in ignore_dirs]
        for name in filenames:
            if patterns:
                if not any(fnmatch.fnmatch(name, pat) for pat in patterns):
                    continue
            path = Path(dirpath) / name
            try:
                if path.stat().st_size > 2 * 1024 * 1024:
                    continue
            except OSError:
                continue
            try:
                with path.open("r", encoding="utf-8", errors="ignore") as f:
                    for lineno, line in enumerate(f, 1):
                        if query in line:
                            rel = path.relative_to(root)
                            snippet = line.strip()
                            matches.append(f"{rel}:{lineno}: {snippet}")
                            if len(matches) >= max_matches:
                                break
                    if len(matches) >= max_matches:
                        break
            except OSError:
                continue
        if len(matches) >= max_matches:
            break

    if not matches:
        logs.append(f"未在 {root} 中搜索到包含「{query}」的内容。")
        return True, {"ok": True, "matches": [], "truncated": False}

    logs.append(f"在 {root} 中找到 {len(matches)} 处匹配（最多展示 {max_matches} 条）。")
    for m in matches:
        logs.append(m)
    truncated = len(matches) >= max_matches
    if truncated:
        logs.append("... 已截断输出，建议缩小范围或使用更精确的关键字。")
    return True, {"ok": True, "matches": matches, "truncated": truncated}


# 工具箱脚本定义: id, name, description, category, params(list of {key, label, type, optional})
TOOLBOX_SCRIPTS = [
    {
        "id": "clear_cache",
        "name": "清理 Lumi 缓存",
        "description": "删除临时工程、缓存文件，释放磁盘空间",
        "category": "维护",
        "params": [],
    },
    {
        "id": "refresh_devices",
        "name": "刷新串口列表",
        "description": "重新扫描 USB 串口设备",
        "category": "设备",
        "params": [],
    },
    {
        "id": "check_mpy",
        "name": "检测 MicroPython",
        "description": "检测指定串口是否为 MicroPython 设备",
        "category": "设备",
        "params": [
            {"key": "port", "label": "串口", "type": "string", "optional": False},
        ],
    },
    {
        "id": "check_pio",
        "name": "检查 PlatformIO",
        "description": "检查 pio 是否安装及版本",
        "category": "环境",
        "params": [],
    },
    {
        "id": "export_code",
        "name": "导出代码到桌面",
        "description": "将当前代码保存为桌面文件",
        "category": "代码",
        "params": [
            {"key": "code", "label": "代码内容", "type": "string", "optional": False},
            {"key": "suffix", "label": "扩展名", "type": "string", "optional": True},
        ],
    },
    {
        "id": "project_search",
        "name": "在项目中搜索文本",
        "description": "在项目根目录下按关键字搜索文件内容（简单版 grep）",
        "category": "代码",
        "params": [
            {"key": "query", "label": "搜索关键字", "type": "string", "optional": False},
            {
                "key": "glob",
                "label": "文件匹配（如 *.py,*.js，可选）",
                "type": "string",
                "optional": True,
            },
        ],
    },
    {
        "id": "show_project_root",
        "name": "显示项目根目录",
        "description": "查看当前 Lumi 使用的项目根路径（AI 编辑、相对路径均基于此）",
        "category": "环境",
        "params": [],
    },
    {
        "id": "open_project_root",
        "name": "在文件管理器中打开项目根",
        "description": "用系统文件管理器打开项目根目录，便于查看/编辑本地文件",
        "category": "环境",
        "params": [],
    },
    {
        "id": "ping_model",
        "name": "检测模型服务",
        "description": "检测您的本地API是否可用及延迟",
        "category": "环境",
        "params": [],
    },
    {
        "id": "check_python",
        "name": "检查 Python 环境",
        "description": "查看本机 Python 版本，确认运行环境",
        "category": "环境",
        "params": [],
    },
    {
        "id": "list_device_files",
        "name": "列出设备上的文件",
        "description": "列出 MicroPython 设备根目录下的文件（需 mpremote）",
        "category": "设备",
        "params": [
            {"key": "port", "label": "串口", "type": "string", "optional": False},
        ],
    },
    {
        "id": "soft_reset_device",
        "name": "软复位 MicroPython 设备",
        "description": "向设备发送软复位，相当于按板子上的复位键",
        "category": "设备",
        "params": [
            {"key": "port", "label": "串口", "type": "string", "optional": False},
        ],
    },
    {
        "id": "mip_install",
        "name": "设备安装 Python 包 (mip)",
        "description": "在 MicroPython 设备上通过 mip 安装包，如 umqtt.simple、aioble",
        "category": "设备",
        "params": [
            {"key": "port", "label": "串口", "type": "string", "optional": False},
            {"key": "package", "label": "包名", "type": "string", "optional": False},
        ],
    },
    {
        "id": "read_device_repl",
        "name": "读取设备 REPL 输出",
        "description": "连接设备并采集一段时间内的输出（内存等状态），用于调试",
        "category": "设备",
        "params": [
            {"key": "port", "label": "串口", "type": "string", "optional": False},
            {"key": "duration", "label": "采集秒数", "type": "string", "optional": True},
        ],
    },
    {
        "id": "run_project_script",
        "name": "在项目根运行 Python 脚本",
        "description": "在 Lumi 项目根目录下执行指定脚本（如测试、工具脚本）",
        "category": "代码",
        "params": [
            {"key": "script_path", "label": "脚本相对路径", "type": "string", "optional": False},
        ],
    },
    {
        "id": "ruff_check",
        "name": "检查项目代码 (ruff)",
        "description": "使用 ruff 对项目根目录进行静态检查",
        "category": "代码",
        "params": [
            {"key": "glob", "label": "文件匹配（如 *.py，可选）", "type": "string", "optional": True},
        ],
    },
    {
        "id": "python_crawler",
        "name": "Python 爬虫 / 下载",
        "description": "输入网址和下载路径，自动爬取并保存到本地（支持网页、图片、文件等）",
        "category": "代码",
        "params": [
            {"key": "url", "label": "要爬取的网址", "type": "string", "optional": False},
            {"key": "download_path", "label": "下载路径（目录或完整路径，留空则保存到桌面）", "type": "string", "optional": True},
        ],
    },
    {
        "id": "pdf_to_word",
        "name": "PDF 转 Word",
        "description": "将 PDF 文件转换为 .docx 文档，输出到指定路径或与 PDF 同目录",
        "category": "文档",
        "params": [
            {"key": "pdf_path", "label": "PDF 文件路径（桌面或项目根下）", "type": "string", "optional": False},
            {"key": "output_path", "label": "输出 .docx 路径（留空则与 PDF 同目录、同名 .docx）", "type": "string", "optional": True},
        ],
    },
    {
        "id": "install_deps",
        "name": "检测并安装依赖",
        "description": "检测本机是否已安装 README 中列出的资源库；若未安装则自动运行 pip 下载安装（requirements.txt + 可选 openai/ruff/platformio），需联网",
        "category": "环境",
        "params": [],
    },
]


def get_toolbox_scripts() -> List[dict]:
    """返回工具箱脚本列表，供前端展示。"""
    return list(TOOLBOX_SCRIPTS)


def run_toolbox_script(
    script_id: str,
    params: dict,
    logs: List[str],
) -> tuple:
    """
    执行指定工具箱脚本。返回 (ok: bool, data: Optional[dict])。
    logs 会被追加执行过程信息。
    """
    def log(msg: str):
        logs.append(msg)
        print(msg)

    if script_id == "clear_cache":
        clear_lumi_cache(logs)
        return True, {"message": "缓存已清理"}

    if script_id == "refresh_devices":
        devices = list_serial_devices()
        guessed = guess_esp8266_port(devices)
        log(f"检测到 {len(devices)} 个串口设备")
        return True, {"devices": devices, "guessed": guessed}

    if script_id == "check_mpy":
        port = (params.get("port") or "").strip()
        ok, msg = probe_micropython(port)
        log(msg)
        return ok, {"ok": ok, "message": msg}

    if script_id == "check_pio":
        ok, msg, version = check_platformio_env()
        log(msg)
        if version:
            log(version)
        return ok, {"ok": ok, "message": msg, "version": version}

    if script_id == "export_code":
        code = params.get("code") or ""
        suffix = (params.get("suffix") or ".py").strip()
        if not suffix.startswith("."):
            suffix = "." + suffix
        if not code:
            log("未提供代码内容，无法导出。")
            return False, {"error": "未提供代码"}
        path = export_code_to_desktop(code, suffix)
        log(f"已导出到: {path}")
        return True, {"path": path}

    if script_id == "project_search":
        query = (params.get("query") or "").strip()
        glob_pattern = (params.get("glob") or "").strip() or None
        ok, data = search_in_project(query, glob_pattern, logs)
        return ok, data

    if script_id == "show_project_root":
        root = get_project_root()
        log(f"项目根目录: {root}")
        return True, {"project_root": root}

    if script_id == "open_project_root":
        try:
            root = open_project_root_in_explorer()
            log(f"已在文件管理器中打开: {root}")
            return True, {"project_root": root}
        except Exception as e:
            log(f"打开失败: {e}")
            return False, {"error": str(e)}

    if script_id == "ping_model":
        result = ping_qwen_model()
        if result.get("ok"):
            log(f"模型服务正常，延迟约 {result.get('latency_ms', 0)} ms")
        else:
            log(f"模型服务异常: {result.get('error', '未知')}")
        return result.get("ok", False), result

    if script_id == "check_python":
        ok, msg, version = check_python_env()
        log(msg)
        if version:
            log(version)
        return ok, {"ok": ok, "message": msg, "version": version}

    if script_id == "install_deps":
        ok, msg = install_missing_dependencies(logs)
        return ok, {"message": msg}

    if script_id == "list_device_files":
        port = (params.get("port") or "").strip()
        ok, out = list_device_files(port)
        log(out)
        return ok, {"ok": ok, "listing": out}

    if script_id == "soft_reset_device":
        port = (params.get("port") or "").strip()
        ok, msg = soft_reset_device(port)
        log(msg)
        return ok, {"ok": ok, "message": msg}

    if script_id == "mip_install":
        port = (params.get("port") or "").strip()
        package = (params.get("package") or "").strip()
        ok, msg = mip_install_on_device(port, package, logs)
        log(msg)
        return ok, {"ok": ok, "message": msg}

    if script_id == "read_device_repl":
        port = (params.get("port") or "").strip()
        duration = (params.get("duration") or "8").strip()
        try:
            duration_sec = int(duration) if duration else 8
        except ValueError:
            duration_sec = 8
        ok, output = read_device_repl(port, duration_sec=duration_sec, logs=logs)
        log(output)
        return ok, {"ok": ok, "output": output}

    if script_id == "run_project_script":
        script_path = (params.get("script_path") or "").strip()
        ok, output = run_project_script(script_path, logs)
        log(output)
        return ok, {"ok": ok, "output": output}

    if script_id == "ruff_check":
        glob_pattern = (params.get("glob") or "").strip() or None
        ok, output = ruff_check_project(glob_pattern, logs)
        log(output)
        return ok, {"ok": ok, "output": output}

    if script_id == "python_crawler":
        url = (params.get("url") or "").strip()
        download_path = (params.get("download_path") or "").strip()
        ok, result = crawl_and_download(url, download_path, logs)
        if ok:
            log(f"已保存到: {result}")
            return True, {"ok": True, "saved_path": result}
        log(result)
        return False, {"ok": False, "error": result}

    if script_id == "pdf_to_word":
        pdf_path = (params.get("pdf_path") or "").strip()
        output_path = (params.get("output_path") or "").strip()
        if not pdf_path:
            log("未填写 PDF 文件路径")
            return False, {"error": "未填写 PDF 文件路径"}
        pdf_path = os.path.normpath(os.path.expanduser(pdf_path))
        if not os.path.isfile(pdf_path):
            log(f"文件不存在: {pdf_path}")
            return False, {"error": f"文件不存在: {pdf_path}"}
        if not pdf_path.lower().endswith(".pdf"):
            log("请指定 .pdf 文件")
            return False, {"error": "请指定 .pdf 文件"}
        if not _is_path_under_allowed_bases(pdf_path):
            log("仅允许转换桌面或项目根下的 PDF 文件")
            return False, {"error": "仅允许转换桌面或项目根下的 PDF 文件"}
        if output_path:
            output_path = os.path.normpath(os.path.expanduser(output_path))
            if not output_path.lower().endswith(".docx"):
                output_path = output_path.rstrip("/") + ".docx"
            out_dir = os.path.dirname(output_path)
            if out_dir and not os.path.isdir(out_dir):
                try:
                    os.makedirs(out_dir, exist_ok=True)
                except OSError as e:
                    log(f"无法创建输出目录: {e}")
                    return False, {"error": str(e)}
            if not _is_path_under_allowed_bases(os.path.abspath(output_path)):
                log("仅允许输出到桌面或项目根下")
                return False, {"error": "仅允许输出到桌面或项目根下"}
        else:
            out_dir = os.path.dirname(pdf_path)
            stem = os.path.splitext(os.path.basename(pdf_path))[0]
            output_path = os.path.join(out_dir, stem + ".docx")
        try:
            from pdf2docx import Converter
            cv = Converter(pdf_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            log(f"已转换: {output_path}")
            return True, {"ok": True, "output_path": output_path}
        except ImportError:
            log("请先安装: pip install pdf2docx")
            return False, {"error": "请先安装: pip install pdf2docx"}
        except Exception as e:
            log(f"转换失败: {e}")
            return False, {"error": str(e)}

    return False, {"error": f"未知脚本: {script_id}"}


# =====================
# 4. 全自动交互流程（命令行 Agent）
# =====================

def interactive_agent():
    # 1. 感知 USB / 串口设备
    devices = list_serial_devices()
    if not devices:
        print("未检测到任何串口设备，请检查 USB 连接。")
        return

    print("检测到以下串口设备：")
    for i, d in enumerate(devices):
        print(f"[{i}] {d['device']}  {d['description']}  ({d.get('manufacturer') or ''})")

    guessed = guess_esp8266_port(devices)
    print()
    print(f"推测 ESP8266 可能在: {guessed}" if guessed else "无法自动推测 ESP8266 串口。")

    idx = input("请选择要操作的设备序号（直接回车使用推测的端口）: ").strip()
    if idx:
        try:
            dev = devices[int(idx)]
            port = dev["device"]
        except Exception:
            print("输入无效，退出。")
            return
    else:
        if not guessed:
            print("没有可用串口，退出。")
            return
        port = guessed

    print(f"将使用串口设备: {port}")

    # 2. 询问用户需求（高层自然语言）
    print()
    print("请用自然语言描述你希望这个 ESP8266 做什么，比如：")
    print(" - 控制一个舵机顺时针转 90 度，然后复位")
    print(" - 连接 Wi-Fi 并定时向服务器上报传感器数据")
    print(" - 控制 LED 呼吸灯并通过串口打印日志")
    user_goal = input("你的需求: ").strip()
    if not user_goal:
        print("未输入需求，退出。")
        return

    # 硬件约束示例，可按自己接线修改
    extra_hardware_context = (
        "\\n硬件约束示例（可根据实际接线修改）：\\n"
        "- 舵机信号线接到 GPIO 5（D1）。\\n"
        "- 使用 machine.PWM 控制舵机，占空比映射到 0~180 度。\\n"
    )

    full_instruction = user_goal + extra_hardware_context

    print("\\n正在调用 Qwen Coder 2.5 生成 MicroPython 代码，请稍候...")
    code = call_qwen_coder(full_instruction)
    print("\\n===== 生成的 main.py 代码预览（前 80 行） =====")
    for i, line in enumerate(code.splitlines()[:80], 1):
        print(f"{i:3}: {line}")
    print("=============================================")
    print()

    confirm = input("是否将此代码覆写到 ESP8266 的 main.py？(y/N): ").strip().lower()
    if confirm != "y":
        print("用户取消写入。")
        return

    # 3. 写临时文件并上传
    tmp_path = write_temp_code(code)
    try:
        flash_micropython_main(port, tmp_path)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    print("已完成本次 Agent 行动。请重启板子或复位观察效果。")


if __name__ == "__main__":
    interactive_agent()

